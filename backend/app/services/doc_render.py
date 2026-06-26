"""Document render engine (cp-0025, restyled cp-0037) — title + sections -> a real,
**designed** PPTX / DOCX / PDF: themed colors, cover banner, section accents, bullet
lists, styled tables, and (PPTX) entrance animations.

Import-guarded + stub-safe (the WORKSPACE-RULE / offline ethos):
  * the format's lib NOT installed (or OMNIVRA_DISABLE_RENDER set) -> returns a stub
    (the caller keeps a markdown deliverable).
  * installed -> a real .pptx (python-pptx) / .docx (python-docx) / .pdf (fpdf2).

Pure/sync + fast (no async needed, unlike video). Never raises: failures return
ok=False; a missing lib returns stub=True. Enable real files with:
    pip install -r requirements-docs.txt
"""
from __future__ import annotations

import hashlib
import math
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from app.core.logging import logger

_FMT_MODULE = {"pptx": "pptx", "docx": "docx", "pdf": "fpdf"}

RGB = tuple[int, int, int]
_WHITE: RGB = (255, 255, 255)


# --- Visual themes ----------------------------------------------------------
@dataclass(frozen=True)
class Theme:
    """A printable, professional palette. ``primary`` fills cover/banners/table headers,
    ``accent`` draws rules + bullets, ``heading`` colors section titles, ``row_alt`` is the
    zebra fill for table data rows, ``ink`` is body text, ``series`` colors chart data."""

    name: str
    primary: RGB
    accent: RGB
    heading: RGB
    ink: RGB
    row_alt: RGB
    series: tuple[RGB, ...] = field(default_factory=tuple)

    def chart_colors(self) -> list[RGB]:
        """Palette for chart series/points — the explicit ``series`` or a primary/accent default."""
        return list(self.series) if self.series else [self.primary, self.accent, self.heading]


THEMES: dict[str, Theme] = {
    "indigo": Theme("indigo", (79, 70, 229), (6, 182, 212), (67, 56, 202), (31, 41, 55), (238, 242, 255),
                    ((79, 70, 229), (6, 182, 212), (168, 85, 247), (244, 114, 182), (52, 211, 153))),
    "emerald": Theme("emerald", (5, 150, 105), (13, 148, 136), (4, 120, 87), (31, 41, 55), (236, 253, 245),
                     ((5, 150, 105), (16, 185, 129), (13, 148, 136), (132, 204, 22), (45, 212, 191))),
    "amber": Theme("amber", (217, 119, 6), (234, 88, 12), (180, 83, 9), (41, 37, 36), (255, 251, 235),
                   ((217, 119, 6), (234, 88, 12), (245, 158, 11), (202, 138, 4), (251, 191, 36))),
    "violet": Theme("violet", (124, 58, 237), (217, 70, 239), (109, 40, 217), (31, 41, 55), (245, 243, 255),
                    ((124, 58, 237), (217, 70, 239), (168, 85, 247), (139, 92, 246), (236, 72, 153))),
    "slate": Theme("slate", (51, 65, 85), (2, 132, 199), (30, 41, 59), (15, 23, 42), (241, 245, 249),
                   ((51, 65, 85), (2, 132, 199), (100, 116, 139), (14, 165, 233), (148, 163, 184))),
    "crimson": Theme("crimson", (190, 18, 60), (244, 63, 94), (159, 18, 57), (35, 20, 25), (255, 241, 242),
                     ((190, 18, 60), (244, 63, 94), (251, 113, 133), (225, 29, 72), (252, 165, 165))),
    "teal": Theme("teal", (13, 148, 136), (6, 182, 212), (15, 118, 110), (19, 42, 46), (240, 253, 250),
                  ((13, 148, 136), (6, 182, 212), (45, 212, 191), (20, 184, 166), (34, 211, 238))),
    "ocean": Theme("ocean", (2, 132, 199), (14, 165, 233), (3, 105, 161), (15, 30, 45), (240, 249, 255),
                   ((2, 132, 199), (14, 165, 233), (56, 189, 248), (59, 130, 246), (99, 102, 241))),
    "sunset": Theme("sunset", (234, 88, 12), (244, 63, 94), (194, 65, 12), (41, 25, 20), (255, 247, 237),
                    ((234, 88, 12), (244, 63, 94), (249, 115, 22), (251, 146, 60), (245, 158, 11))),
    "forest": Theme("forest", (22, 101, 52), (101, 163, 13), (20, 83, 45), (20, 35, 25), (240, 253, 244),
                    ((22, 101, 52), (101, 163, 13), (34, 197, 94), (132, 204, 22), (74, 222, 128))),
    "midnight": Theme("midnight", (30, 41, 59), (99, 102, 241), (15, 23, 42), (15, 23, 42), (241, 245, 249),
                      ((30, 41, 59), (99, 102, 241), (129, 140, 248), (56, 189, 248), (167, 139, 250))),
    "rose": Theme("rose", (190, 24, 93), (236, 72, 153), (157, 23, 77), (40, 20, 30), (253, 242, 248),
                  ((190, 24, 93), (236, 72, 153), (244, 114, 182), (219, 39, 119), (251, 182, 206))),
}
_DEFAULT_THEME = "indigo"

# Layout/style variants — vary the deck's look (cover + section header treatment) so documents
# don't all look identical. Picked deterministically from the title (stable on regenerate, varied
# across topics) by :func:`pick_style`.
STYLES = ("band", "sidebar", "minimal")


def resolve_theme(name: str | None) -> Theme:
    """Map a theme name to its palette; 'auto'/unknown/None -> the default (indigo)."""
    return THEMES.get((name or "").strip().lower(), THEMES[_DEFAULT_THEME])


def pick_style(seed: str) -> str:
    """Deterministically choose a layout style from ``seed`` (the title) — stable per document,
    varied across documents (the 'different style each time' the user asked for)."""
    h = int(hashlib.sha1((seed or "doc").encode("utf-8", "replace")).hexdigest(), 16)
    return STYLES[h % len(STYLES)]


# --- Per-genre STRUCTURE specs ----------------------------------------------
# Each writing style is a recognizable document GENRE. The StructureSpec captures its visual
# STRUCTURE only — cover, columns, numbering, page frame, heading treatment, dividers, containers,
# callouts, tables, image role. Color (`theme`) and typeface (`font`) are chosen INDEPENDENTLY by
# the user, so two genres at the same theme+font must still look nothing alike — all the difference
# lives here. Defaults describe a clean corporate report; each genre overrides what makes it distinct.
@dataclass(frozen=True)
class StructureSpec:
    genre: str = "Report"
    # page frame
    page_border: str = "none"            # none | full | top_bottom
    margin: str = "normal"               # tight | normal | wide
    columns: int = 1                     # 1 | 2  (prose flows in 2 cols; blocks stay full-width)
    header: str = "none"                 # none | running_head | masthead | letterhead
    footer: str = "none"                 # none | page_num | page_x_of_y | brand
    # cover / title
    cover: str = "band"                  # inline|band|fullbleed_image|centered|letterhead|hero_number|magazine|plain
    cover_page: bool = False             # dedicated first page (PDF/DOCX); PPTX always has a cover slide
    drop_cap: bool = False               # large initial on the first section's body
    # headings
    numbering: str = "none"              # none | decimal (1, 1.1) | legal (1.1.1) | section_word ("Section 1")
    heading_style: str = "rule_under"    # rule_under|left_bar|boxed|kicker|allcaps|bottom_border|plain
    heading_align: str = "left"          # left | center
    # body & lists
    justify: bool = False
    line_spacing: float = 1.3
    body_size: int = 11                  # base body pt (PDF); DOCX/PPTX scale from this
    bullet: str = "dot"                  # dot|dash|triangle|arrow|diamond|check|numbered|none
    # sections
    divider: str = "hairline"            # none|hairline|thick|shaded|ornament|new_page
    container: str = "flow"              # flow|card|shaded|bordered
    # callouts
    callout: str = "tint"                # none|tint|quote_bar|warning|pull_quote|speech
    # tables
    table_style: str = "zebra"           # grid|zebra|minimal|boxed_header|borderless
    # images (FLUX)
    image_policy: str = "optional"       # none|optional|inline|hero|fullbleed
    image_max: int = 1                   # hard cap on generated images per document
    # extras
    toc: bool = False                    # render a table of contents
    cta: bool = False                    # render a closing call-to-action banner


STRUCTURE: dict[str, StructureSpec] = {
    "professional": StructureSpec(
        genre="Corporate report", cover="band", header="running_head", footer="page_num",
        heading_style="rule_under", line_spacing=1.3, bullet="dot", divider="hairline",
        container="flow", callout="tint", table_style="zebra", image_policy="optional", image_max=1),
    "casual": StructureSpec(
        genre="Blog / zine", cover="plain", margin="wide", heading_style="plain", line_spacing=1.5,
        body_size=12, bullet="triangle", divider="none", container="flow", callout="tint",
        table_style="minimal", image_policy="inline", image_max=2),
    "academic": StructureSpec(
        genre="Journal article", cover="centered", cover_page=True, columns=2, header="running_head",
        footer="page_num", numbering="decimal", heading_style="plain", justify=True, line_spacing=1.5,
        body_size=10, bullet="dash", divider="none", container="flow", callout="quote_bar",
        table_style="minimal", image_policy="none", image_max=0),
    "formal": StructureSpec(
        genre="Official memo", cover="letterhead", header="letterhead", footer="page_x_of_y",
        numbering="decimal", heading_style="allcaps", justify=True, line_spacing=1.4, bullet="dash",
        divider="hairline", container="flow", callout="none", table_style="grid",
        image_policy="none", image_max=0),
    "informal": StructureSpec(
        genre="Notebook", cover="plain", margin="wide", heading_style="left_bar", line_spacing=1.5,
        body_size=12, bullet="dash", divider="none", container="flow", callout="tint",
        table_style="borderless", image_policy="optional", image_max=1),
    "conversational": StructureSpec(
        genre="Q & A", cover="plain", heading_style="plain", line_spacing=1.5, body_size=12,
        bullet="arrow", divider="hairline", container="flow", callout="speech", table_style="minimal",
        image_policy="optional", image_max=1),
    "technical": StructureSpec(
        genre="Manual / spec", cover="band", header="running_head", footer="page_num",
        numbering="decimal", heading_style="left_bar", line_spacing=1.25, body_size=10,
        bullet="numbered", divider="hairline", container="bordered", callout="warning",
        table_style="grid", image_policy="optional", image_max=1, toc=True),
    "business": StructureSpec(
        genre="Executive briefing", cover="hero_number", header="running_head", footer="page_num",
        heading_style="bottom_border", line_spacing=1.3, bullet="dot", divider="shaded",
        container="shaded", callout="tint", table_style="zebra", image_policy="optional", image_max=1),
    "creative": StructureSpec(
        genre="Magazine feature", cover="magazine", margin="wide", columns=2, footer="page_num",
        drop_cap=True, heading_style="kicker", justify=True, line_spacing=1.4, bullet="diamond",
        divider="ornament", container="flow", callout="pull_quote", table_style="minimal",
        image_policy="hero", image_max=2),
    "simple": StructureSpec(
        genre="Minimalist", cover="centered", margin="wide", heading_style="plain", line_spacing=1.6,
        body_size=13, bullet="dash", divider="none", container="flow", callout="none",
        table_style="borderless", image_policy="none", image_max=0),
    "complex": StructureSpec(
        genre="Whitepaper", cover="centered", cover_page=True, margin="tight", columns=2,
        header="running_head", footer="page_x_of_y", numbering="legal", heading_style="plain",
        justify=True, line_spacing=1.2, body_size=10, bullet="numbered", divider="hairline",
        container="flow", callout="quote_bar", table_style="grid", image_policy="optional",
        image_max=1, toc=True),
    "concise": StructureSpec(
        genre="One-pager brief", cover="inline", page_border="top_bottom", heading_style="rule_under",
        line_spacing=1.15, bullet="dot", divider="hairline", container="flow", callout="tint",
        table_style="minimal", image_policy="none", image_max=0),
    "detailed": StructureSpec(
        genre="Comprehensive doc", cover="band", cover_page=True, header="running_head",
        footer="page_x_of_y", numbering="decimal", heading_style="rule_under", justify=True,
        line_spacing=1.4, bullet="dot", divider="hairline", container="flow", callout="tint",
        table_style="zebra", image_policy="optional", image_max=1, toc=True),
    "persuasive": StructureSpec(
        genre="Proposal / pitch", cover="band", footer="page_num", heading_style="bottom_border",
        line_spacing=1.35, body_size=12, bullet="arrow", divider="thick", container="flow",
        callout="pull_quote", table_style="zebra", image_policy="hero", image_max=2, cta=True),
    "informative": StructureSpec(
        genre="Newsletter", cover="masthead", columns=2, header="masthead", footer="page_num",
        heading_style="kicker", justify=True, line_spacing=1.35, body_size=10, bullet="dot",
        divider="hairline", container="flow", callout="tint", table_style="minimal",
        image_policy="inline", image_max=2),
    "neutral": StructureSpec(
        genre="Standards doc", cover="centered", page_border="full", header="running_head",
        footer="page_x_of_y", numbering="decimal", heading_style="plain", justify=True,
        line_spacing=1.4, bullet="dash", divider="hairline", container="flow", callout="none",
        table_style="grid", image_policy="none", image_max=0),
    "friendly": StructureSpec(
        genre="Brochure", cover="band", margin="wide", heading_align="center", heading_style="boxed",
        line_spacing=1.45, body_size=12, bullet="check", divider="none", container="card",
        callout="tint", table_style="borderless", image_policy="inline", image_max=2),
    "seo-friendly": StructureSpec(
        genre="Web article", cover="plain", heading_style="bottom_border", line_spacing=1.5,
        body_size=12, bullet="check", divider="hairline", container="flow", callout="tint",
        table_style="minimal", image_policy="inline", image_max=1),
    "marketing": StructureSpec(
        genre="Landing page", cover="fullbleed_image", margin="wide", heading_align="center",
        heading_style="kicker", line_spacing=1.4, body_size=12, bullet="check", divider="thick",
        container="card", callout="pull_quote", table_style="zebra", image_policy="fullbleed",
        image_max=2, cta=True),
    "legal": StructureSpec(
        genre="Contract / brief", cover="letterhead", cover_page=True, page_border="full",
        header="running_head", footer="page_x_of_y", numbering="legal", heading_style="allcaps",
        justify=True, line_spacing=1.5, bullet="numbered", divider="none", container="flow",
        callout="none", table_style="grid", image_policy="none", image_max=0, toc=True),
}
_DEFAULT_STRUCT = STRUCTURE["professional"]


def structure_spec(style: str | None) -> StructureSpec:
    """Map a writing style/genre to its visual StructureSpec; unknown -> the professional default."""
    return STRUCTURE.get((style or "").strip().lower(), _DEFAULT_STRUCT)


# Font family -> concrete font names per format. PDF (fpdf2) has only core families; DOCX/PPTX
# can name any installed font (these are standard Office fonts). Chosen by the USER, not the genre.
_PDF_FONT = {"serif": "Times", "sans": "Helvetica", "mono": "Courier"}
_OFFICE_FONT = {"serif": "Georgia", "sans": "Calibri", "mono": "Consolas"}
# Bullet glyphs: PDF uses latin-1-safe marks; DOCX/PPTX can use real Unicode glyphs.
_BULLET_PDF = {"dot": "\xb7", "dash": "-", "triangle": ">", "arrow": ">", "diamond": "*", "check": "*"}
_BULLET_UNI = {"dot": "•", "dash": "–", "triangle": "▸", "arrow": "→", "diamond": "◆", "check": "✓"}
_MARGIN_MM = {"tight": 12.0, "normal": 16.0, "wide": 22.0}


def pdf_font(font: str | None) -> str:
    """Resolve the user's font family to an fpdf2 core family name (Times/Helvetica/Courier)."""
    return _PDF_FONT.get((font or "").strip().lower(), "Helvetica")


def office_font(font: str | None) -> str:
    """Resolve the user's font family to an Office font name (Georgia/Calibri/Consolas)."""
    return _OFFICE_FONT.get((font or "").strip().lower(), "Calibri")


def _hex(rgb: RGB) -> str:
    return "%02X%02X%02X" % rgb


def doc_available(fmt: str) -> bool:
    """True when the render lib for ``fmt`` is importable (and rendering isn't disabled)."""
    if os.environ.get("OMNIVRA_DISABLE_RENDER"):
        return False
    mod = _FMT_MODULE.get(fmt)
    if not mod:
        return False
    try:
        __import__(mod)
        return True
    except Exception:  # noqa: BLE001
        return False


# fpdf2 core fonts are latin-1, so common "smart" Unicode punctuation (em dash, curly quotes,
# ellipsis, arrows, bullets) would otherwise become "?". Transliterate to clean ASCII first.
_LATIN1_TRANSLIT = {
    "‘": "'", "’": "'", "‚": "'", "‛": "'",       # single quotes
    "“": '"', "”": '"', "„": '"', "«": '"', "»": '"',  # double quotes
    "–": "-", "—": "-", "―": "-", "−": "-",       # dashes / minus
    "…": "...",                                                    # ellipsis
    "•": "-", "·": "-", "●": "-", "▪": "-", "⁃": "-",  # bullets
    "→": "->", "←": "<-", "⇒": "=>", "↔": "<->",  # arrows
    "✓": "*", "✔": "*", "✗": "x",                        # check / cross marks
    " ": " ", " ": " ", " ": " ", "​": "",          # spaces
    "™": "(TM)", "®": "(R)", "©": "(C)", "€": "EUR",
}
_LATIN1_TABLE = {ord(k): v for k, v in _LATIN1_TRANSLIT.items()}


def _latin1(text: str) -> str:
    """Make ``text`` safe for fpdf2 core (latin-1) fonts: transliterate common Unicode punctuation
    to ASCII first (so it reads cleanly, not as '?'), then drop anything still unencodable."""
    return (text or "").translate(_LATIN1_TABLE).encode("latin-1", "replace").decode("latin-1")


def _norm_table(table: Any) -> tuple[list[str], list[list[str]]] | None:
    """Normalize a section's ``table`` (dict | None) into (headers, rows) of strings, or None."""
    if not isinstance(table, dict):
        return None
    headers = [str(h) for h in (table.get("headers") or [])]
    rows = [[str(c) for c in row] for row in (table.get("rows") or []) if isinstance(row, (list, tuple))]
    if not headers and not rows:
        return None
    if not headers and rows:  # rows without an explicit header — synthesize blank headers
        headers = [""] * max((len(r) for r in rows), default=0)
    ncols = len(headers)
    if ncols == 0:  # degenerate (no headers AND only empty rows) -> treat as 'no table' everywhere
        return None
    rows = [(r + [""] * ncols)[:ncols] for r in rows]  # pad/clip every row to header width
    return headers, rows


_CHART_TYPES = {"column", "bar", "line", "pie", "area"}


def _num(v: Any) -> float:
    """Coerce a possibly-stringy value to a FINITE float (strip %, $, commas); non-numeric or
    non-finite (inf/nan) -> 0.0 so it never breaks the chart's embedded worksheet writer."""
    if isinstance(v, (int, float)):
        f = float(v)
        return f if math.isfinite(f) else 0.0
    s = str(v).strip().replace(",", "").replace("$", "").replace("%", "")
    try:
        f = float(s)
    except (TypeError, ValueError):
        return 0.0
    return f if math.isfinite(f) else 0.0


def _norm_chart(chart: Any) -> tuple[str, str, list[str], list[tuple[str, list[float]]]] | None:
    """Normalize a section's ``chart`` (dict | None) into (type, title, categories, series), or None
    when there's no usable numeric data. Categories/series are padded so every series aligns."""
    if not isinstance(chart, dict):
        return None
    ctype = str(chart.get("type", "column")).strip().lower()
    if ctype not in _CHART_TYPES:
        ctype = "column"
    title = str(chart.get("title", "")).strip()
    categories = [str(c).strip() for c in (chart.get("categories") or [])]
    series: list[tuple[str, list[float]]] = []
    for s in (chart.get("series") or []):
        if not isinstance(s, dict):
            continue
        values = [_num(v) for v in (s.get("values") or [])]
        if values:
            series.append((str(s.get("name", "")).strip(), values))
    if not series:
        return None
    width = max(len(categories), max(len(v) for _, v in series))
    if width == 0:
        return None
    if len(categories) < width:  # synthesize missing category labels
        categories += [str(i + 1) for i in range(len(categories), width)]
    categories = categories[:width]
    series = [(name, (vals + [0.0] * width)[:width]) for name, vals in series]
    if ctype == "pie" and len(series) > 1:  # a pie is single-series — keep only the first so
        series = series[:1]                  # PPTX and the DOCX/PDF data table agree
    return ctype, title, categories, series


def _chart_to_table(chart_tuple: tuple[str, str, list[str], list[tuple[str, list[float]]]]) -> tuple[str, tuple[list[str], list[list[str]]]]:
    """Render a normalized chart as a (caption, table) pair for DOCX/PDF (which don't draw charts).
    Caption is the chart title; the table is Category + one column per series."""
    ctype, title, categories, series = chart_tuple
    headers = ["Category"] + [name or f"Series {i+1}" for i, (name, _) in enumerate(series)]

    def _fmt(x: float) -> str:
        return str(int(x)) if float(x).is_integer() else f"{x:g}"

    rows = [[cat] + [_fmt(vals[i]) for _, vals in series] for i, cat in enumerate(categories)]
    return (title or f"{ctype.title()} chart"), (headers, rows)


# ===========================================================================
# PPTX
# ===========================================================================
def _render_pptx(
    title: str, subtitle: str, sections: list[dict[str, Any]], out: Path, theme: Theme,
    *, font: str = "sans", spec: "StructureSpec | None" = None, asset_root: Path | None = None,
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt

    sp = spec or _DEFAULT_STRUCT
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height
    primary, accent = RGBColor(*theme.primary), RGBColor(*theme.accent)
    ink, heading_c = RGBColor(*theme.ink), RGBColor(*theme.heading)
    white, light, muted = RGBColor(*_WHITE), RGBColor(248, 250, 252), RGBColor(140, 140, 150)
    font_name = office_font(font)
    glyph = _BULLET_UNI.get(sp.bullet, "•")
    hs = sp.heading_style

    def rect(slide, left, top, width, height, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tb, tf

    def put(tf, text, size, color, *, bold=False, italic=False, para=None, align=None):
        p = tf.paragraphs[0] if para is None else para
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size, r.font.bold, r.font.italic, r.font.color.rgb = Pt(size), bold, italic, color
        r.font.name = font_name  # user typeface (serif/sans/mono)
        return r

    def cover_image() -> Path | None:
        if sp.cover not in ("fullbleed_image", "magazine"):
            return None
        for s in sections:
            ap = _img_abs(s, asset_root)
            if ap is not None:
                return ap
        return None

    cov_img = cover_image()
    consumed_id = None
    if cov_img is not None:
        for s in sections:
            if _img_abs(s, asset_root) == cov_img:
                consumed_id = id(s); break

    # ---- Cover slide (per genre) ----
    cover = prs.slides.add_slide(blank)
    cv = sp.cover
    if cv == "fullbleed_image" and cov_img is not None:
        try:
            cover.shapes.add_picture(str(cov_img), 0, 0, width=W, height=H)
        except Exception:  # noqa: BLE001
            rect(cover, 0, 0, W, H, primary)
        rect(cover, 0, Inches(5.2), W, Inches(2.3), ink)
        _, tf = textbox(cover, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.4), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 50, white, bold=True, align=PP_ALIGN.CENTER)
    elif cv == "magazine":
        if cov_img is not None:
            try:
                cover.shapes.add_picture(str(cov_img), 0, 0, width=W, height=Inches(4.6))
            except Exception:  # noqa: BLE001
                rect(cover, 0, 0, W, Inches(4.6), primary)
        else:
            rect(cover, 0, 0, W, Inches(4.6), primary)
        _, tf = textbox(cover, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.4))
        put(tf, (sp.genre or "Feature").upper(), 16, accent, bold=True)
        _, tf = textbox(cover, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6))
        put(tf, title or "Untitled", 46, primary, bold=True)
    elif cv == "centered":
        rect(cover, 0, 0, W, H, light)
        _, tf = textbox(cover, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.8), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 54, primary, bold=True, align=PP_ALIGN.CENTER)
        rect(cover, Inches(5.4), Inches(4.5), Inches(2.5), Inches(0.07), accent)
        _, tf = textbox(cover, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.0))
        put(tf, subtitle or "Omnivra AI Company OS", 24, ink, align=PP_ALIGN.CENTER)
    elif cv == "letterhead":
        rect(cover, 0, 0, W, H, light)
        rect(cover, 0, 0, W, Inches(0.18), primary)
        _, tf = textbox(cover, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.5))
        put(tf, "OMNIVRA AI COMPANY OS", 14, primary, bold=True)
        rect(cover, Inches(0.9), Inches(3.4), Inches(2.4), Inches(0.09), accent)
        _, tf = textbox(cover, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.7), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 44, ink, bold=True)
        _, tf = textbox(cover, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.0))
        put(tf, subtitle or "", 22, muted, italic=True)
    elif cv == "masthead":
        rect(cover, 0, 0, W, H, light)
        rect(cover, 0, Inches(2.5), W, Inches(0.06), ink)
        _, tf = textbox(cover, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.6), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 56, primary, bold=True, align=PP_ALIGN.CENTER)
        rect(cover, 0, Inches(4.3), W, Inches(0.06), ink)
        _, tf = textbox(cover, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5))
        put(tf, subtitle or sp.genre, 18, muted, italic=True, align=PP_ALIGN.CENTER)
    elif cv == "hero_number":
        rect(cover, 0, 0, W, H, light)
        rect(cover, 0, 0, Inches(0.6), H, primary)
        _, tf = textbox(cover, Inches(1.2), Inches(2.4), Inches(11.0), Inches(1.9), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 52, primary, bold=True)
        _, tf = textbox(cover, Inches(1.2), Inches(4.4), Inches(11.0), Inches(1.0))
        put(tf, subtitle or "Omnivra AI Company OS", 24, ink)
    elif cv == "plain":
        rect(cover, 0, 0, W, H, light)
        _, tf = textbox(cover, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.0), MSO_ANCHOR.MIDDLE)
        put(tf, title or "Untitled", 60, primary, bold=True)
        _, tf = textbox(cover, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.0))
        put(tf, subtitle or "", 26, muted)
    else:  # band / inline -> full-bleed primary cover
        rect(cover, 0, 0, W, H, primary)
        rect(cover, Inches(0.9), Inches(4.2), Inches(2.6), Inches(0.13), accent)
        _, tf = textbox(cover, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.7), MSO_ANCHOR.BOTTOM)
        put(tf, title or "Untitled", 56, white, bold=True)
        _, tf = textbox(cover, Inches(0.9), Inches(4.45), Inches(11.5), Inches(1.0))
        put(tf, subtitle or "Generated by Omnivra AI Company OS", 26, white)

    def section_header(slide, heading: str, idx: int) -> tuple[float, float, float]:
        """Draw the per-genre section header; return (content_left, content_width, content_top) in in."""
        num = _sec_number(sp.numbering, idx)
        label = f"{num}  {heading}" if num else heading
        if hs == "left_bar":
            rect(slide, 0, 0, Inches(0.45), H, primary)
            _, tf = textbox(slide, Inches(0.85), Inches(0.4), Inches(12.0), Inches(0.95), MSO_ANCHOR.MIDDLE)
            put(tf, label, 34, heading_c, bold=True)
            rect(slide, Inches(0.85), Inches(1.42), Inches(11.9), Inches(0.05), accent)
            return 0.85, 11.9, 1.72
        if hs == "kicker":
            _, tf = textbox(slide, Inches(0.7), Inches(0.35), Inches(11.0), Inches(0.4))
            put(tf, (sp.genre or "Section").upper(), 13, accent, bold=True)
            _, tf = textbox(slide, Inches(0.7), Inches(0.72), Inches(12.0), Inches(0.9), MSO_ANCHOR.MIDDLE)
            put(tf, label, 34, heading_c, bold=True)
            rect(slide, Inches(0.7), Inches(1.62), Inches(4.6), Inches(0.05), accent)
            return 0.7, 11.95, 1.9
        if hs == "boxed":
            rect(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.95), RGBColor(*theme.row_alt))
            rect(slide, Inches(0.6), Inches(0.35), Inches(0.08), Inches(0.95), accent)
            _, tf = textbox(slide, Inches(0.9), Inches(0.35), Inches(11.6), Inches(0.95), MSO_ANCHOR.MIDDLE)
            put(tf, label, 32, heading_c, bold=True,
                align=PP_ALIGN.CENTER if sp.heading_align == "center" else PP_ALIGN.LEFT)
            return 0.7, 11.95, 1.6
        if hs in ("bottom_border", "plain") or hs == "allcaps" and sp.cover in ("plain", "centered"):
            _, tf = textbox(slide, Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.95), MSO_ANCHOR.MIDDLE)
            put(tf, label.upper() if hs == "allcaps" else label, 36, heading_c, bold=True,
                align=PP_ALIGN.CENTER if sp.heading_align == "center" else PP_ALIGN.LEFT)
            rect(slide, Inches(0.7), Inches(1.45), Inches(3.0), Inches(0.08), accent)
            return 0.7, 11.95, 1.82
        # band (default: rule_under / allcaps on band)
        rect(slide, 0, 0, W, Inches(1.2), primary)
        rect(slide, 0, Inches(1.2), W, Inches(0.07), accent)
        _, tf = textbox(slide, Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.85), MSO_ANCHOR.MIDDLE)
        put(tf, label.upper() if hs == "allcaps" else label, 36, white, bold=True)
        return 0.6, 12.1, 1.65

    FLOOR = 6.9

    def draw_table(slide, headers, rows, cl, cw, top, max_h) -> tuple[list[int], float]:
        ncols, per_row = len(headers), 0.52
        fit = max(1, int(max_h / per_row) - 1)
        truncated = len(rows) > fit
        if truncated:
            fit = max(1, int((max_h - 0.32) / per_row) - 1)
        shown = rows[:fit]
        nrows = len(shown) + 1
        theight = min(max_h - (0.32 if truncated else 0.0), per_row * nrows)
        gf = slide.shapes.add_table(nrows, ncols, Inches(cl), Inches(top), Inches(cw), Inches(theight))
        tbl = gf.table
        col_w = int(Inches(cw) / ncols)
        minimal = sp.table_style in ("minimal", "borderless")
        for c in range(ncols):
            tbl.columns[c].width = col_w
        for c, htext in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = white if minimal else primary
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = htext
            pf = cell.text_frame.paragraphs[0]
            pf.alignment = PP_ALIGN.LEFT
            pf.font.bold, pf.font.size, pf.font.color.rgb, pf.font.name = (
                True, Pt(18), heading_c if minimal else white, font_name)
        alt = RGBColor(*theme.row_alt)
        zebra = sp.table_style in ("zebra", "boxed_header")
        for ri, row in enumerate(shown, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(ri, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt if (zebra and ri % 2 == 0) else white
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text = val
                pf = cell.text_frame.paragraphs[0]
                pf.font.size, pf.font.color.rgb, pf.font.name = Pt(16), ink, font_name
        ids, used = [gf.shape_id], theight
        if truncated:
            tb, tf = textbox(slide, Inches(cl), Inches(top + theight + 0.02), Inches(cw), Inches(0.28))
            put(tf, f"+{len(rows) - len(shown)} more rows", 12, muted)
            ids.append(tb.shape_id)
            used += 0.32
        return ids, used

    # ---- Section slides ----
    for idx, sec in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)
        anim_targets: list[int] = []
        cl, cw, ctop = section_header(slide, sec.get("heading", ""), idx)
        cur = ctop
        body = (sec.get("body") or "").strip()
        bullets = [str(b).strip() for b in (sec.get("bullets") or []) if str(b).strip()]
        chart = _norm_chart(sec.get("chart"))
        table = _norm_table(sec.get("table"))
        has_visual = bool(chart or table)
        pull = sp.callout == "pull_quote" and body and len(body) < 240 and not bullets

        if body:
            room = max(0.6, FLOOR - cur)
            bx = max(0.7, min(2.6, room * (0.42 if has_visual else (0.72 if bullets else 1.0))))
            tb, tf = textbox(slide, Inches(cl), Inches(cur), Inches(cw), Inches(bx),
                             MSO_ANCHOR.MIDDLE if pull else MSO_ANCHOR.TOP)
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if pull:
                put(tf, f"“{body}”", 30, primary, bold=True, italic=True, align=PP_ALIGN.CENTER)
            else:
                put(tf, body, 26, ink)
            anim_targets.append(tb.shape_id)
            cur += bx + 0.12

        if bullets:
            reserve = 2.2 if has_visual else 0.0
            avail = max(0.6, FLOOR - cur - reserve)
            per = 0.62
            max_b = max(1, int(avail / per))
            shown = bullets[:max_b]
            more = len(bullets) - len(shown)
            box_h = min(avail, len(shown) * per + (0.32 if more else 0.0))
            tb, tf = textbox(slide, Inches(cl + 0.1), Inches(cur), Inches(cw - 0.1), Inches(box_h))
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for i, b in enumerate(shown):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                mark = f"{i + 1}.  " if sp.bullet == "numbered" else (f"{glyph}  " if sp.bullet != "none" else "")
                if mark:
                    dot = para.add_run()
                    dot.text = mark
                    dot.font.size, dot.font.color.rgb, dot.font.name = Pt(24), accent, font_name
                run = para.add_run()
                run.text = b
                run.font.size, run.font.color.rgb, run.font.name = Pt(26), ink, font_name
                para.space_after = Pt(8)
            if more:
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = f"(+{more} more)"
                r.font.size, r.font.italic, r.font.color.rgb, r.font.name = Pt(14), True, muted, font_name
            anim_targets.append(tb.shape_id)
            cur += box_h + 0.15

        if chart:
            room = FLOOR - cur
            if room >= 1.7:
                try:
                    ch_h = min(4.2, room)
                    gf = _pptx_chart(slide, chart, Inches(cl), Inches(cur), Inches(cw), Inches(ch_h), theme)
                    anim_targets.append(gf.shape_id)
                    cur += ch_h + 0.12
                except Exception as exc:  # noqa: BLE001
                    logger.warning("pptx chart fell back to a table: {}", repr(exc))
                    if table is None:
                        table = _chart_to_table(chart)[1]
            elif table is None:
                table = _chart_to_table(chart)[1]

        if table:
            room = FLOOR - cur
            if room >= 0.7:
                ids, used = draw_table(slide, table[0], table[1], cl, cw, cur, room)
                anim_targets.extend(ids)
                cur += used + 0.1

        rect(slide, 0, Inches(7.35), W, Inches(0.05), accent)
        if sp.footer != "none":
            _, tf = textbox(slide, Inches(12.2), Inches(6.95), Inches(0.9), Inches(0.4), MSO_ANCHOR.MIDDLE)
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            put(tf, str(idx), 11, muted)
        try:
            _pptx_add_fade_in(slide, anim_targets)
        except Exception as exc:  # noqa: BLE001
            logger.warning("pptx animation injection skipped: {}", repr(exc))

        # Dedicated image slide for a section illustration (kept off the busy text slide).
        ap = _img_abs(sec, asset_root)
        if ap is not None and id(sec) != consumed_id and sp.image_policy != "none":
            try:
                islide = prs.slides.add_slide(blank)
                rect(islide, 0, 0, W, Inches(1.0), primary)
                _, tf = textbox(islide, Inches(0.6), Inches(0.1), Inches(12.1), Inches(0.8), MSO_ANCHOR.MIDDLE)
                put(tf, sec.get("heading", ""), 28, white, bold=True)
                islide.shapes.add_picture(str(ap), Inches(1.4), Inches(1.3), height=Inches(5.6))
            except Exception as exc:  # noqa: BLE001
                logger.warning("pptx image slide skipped: {}", repr(exc))

    prs.save(str(out))


def _pptx_chart(slide, chart_tuple, x, y, cx, cy, theme: Theme):
    """Draw a native PPTX chart (column/bar/line/pie/area) themed with the palette. Returns the
    GraphicFrame. Raises on bad data so the caller can fall back to a data table."""
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Pt

    ctype, title, categories, series = chart_tuple
    type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }
    xltype = type_map.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)

    cd = CategoryChartData()
    cd.categories = categories
    if ctype == "pie":  # pie shows a single series
        name, values = series[0]
        cd.add_series(name or title or "Series", values)
    else:
        for name, values in series:
            cd.add_series(name or "Series", values)

    gframe = slide.shapes.add_chart(xltype, x, y, cx, cy, cd)
    chart = gframe.chart
    palette = theme.chart_colors()

    chart.has_title = bool(title)
    if title:
        tf = chart.chart_title.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True

    multi = len(series) > 1 or ctype == "pie"
    chart.has_legend = multi
    if multi:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(14)

    plot = chart.plots[0]
    if ctype == "pie":
        plot.has_data_labels = True
        # show_percentage makes PowerPoint compute each slice's SHARE of the whole; without it,
        # numFmt '0%' would format the raw value (e.g. 70 -> '7000%').
        plot.data_labels.show_percentage = True
        plot.data_labels.show_value = False
        plot.data_labels.number_format = "0%"
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.size = Pt(14)
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(*palette[i % len(palette)])
    else:
        for i, s in enumerate(chart.series):
            col = RGBColor(*palette[i % len(palette)])
            if ctype == "line":
                s.format.line.color.rgb = col
            else:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = col
        try:  # bigger axis tick labels for readability
            chart.category_axis.tick_labels.font.size = Pt(14)
            chart.value_axis.tick_labels.font.size = Pt(14)
        except Exception:  # noqa: BLE001 - some chart types lack one of the axes
            pass
    return gframe


def _pptx_add_fade_in(slide, shape_ids: list[int]) -> None:
    """Inject a sequential Fade ('Start After Previous') entrance animation onto ``shape_ids``.

    python-pptx has no animation API, so we append a standard ``<p:timing>`` tree to the slide
    XML. presetID=10/presetClass=entr/filter=fade is PowerPoint's Fade entrance. Each target gets
    its own ``<p:par>`` that auto-plays after the prior one (small delay), so the slide animates on
    entry without a click. Caller wraps this in try/except — a bad tree must never break the file.
    """
    if not shape_ids:
        return
    from pptx.oxml import parse_xml

    # Frozen OOXML PresentationML namespace (ECMA-376). pptx.oxml.ns.nsmap is a function,
    # not a dict, so we declare the prefix ourselves on the <p:timing> root.
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    cid = 100  # unique ids across the tree; start high to avoid any collision

    def nid() -> int:
        nonlocal cid
        cid += 1
        return cid

    blocks = []
    for idx, spid in enumerate(shape_ids):
        delay = "0" if idx == 0 else "300"  # first on the slide's click, rest 0.3s after the previous
        # Canonical PowerPoint sequencing: first effect starts on click, the rest "after previous".
        node = "clickEffect" if idx == 0 else "afterEffect"
        blocks.append(
            f'<p:par><p:cTn id="{nid()}" fill="hold"><p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
            f'<p:childTnLst><p:par><p:cTn id="{nid()}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par><p:cTn id="{nid()}" presetID="10" presetClass="entr" presetSubtype="0" '
            f'fill="hold" grpId="0" nodeType="{node}"><p:childTnLst>'
            f'<p:set><p:cBhvr><p:cTn id="{nid()}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            f'<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{nid()}" dur="500"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
            f'</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
        )
    timing_xml = (
        f'<p:timing xmlns:p="{p_ns}">'
        f'<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'{"".join(blocks)}'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    slide._element.append(parse_xml(timing_xml))  # CT_Slide: timing is the last (optional) child


# ===========================================================================
# DOCX
# ===========================================================================
def _render_docx(
    title: str, subtitle: str, sections: list[dict[str, Any]], out: Path, theme: Theme,
    *, font: str = "sans", spec: "StructureSpec | None" = None, asset_root: Path | None = None,
) -> None:
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor, Inches

    sp = spec or _DEFAULT_STRUCT
    ofont = office_font(font)
    primary, accent, heading_c, ink = (
        RGBColor(*theme.primary), RGBColor(*theme.accent), RGBColor(*theme.heading), RGBColor(*theme.ink),
    )
    muted = RGBColor(120, 120, 128)
    align_body = WD_ALIGN_PARAGRAPH.JUSTIFY if sp.justify else WD_ALIGN_PARAGRAPH.LEFT
    halign = WD_ALIGN_PARAGRAPH.CENTER if sp.heading_align == "center" else WD_ALIGN_PARAGRAPH.LEFT
    glyph = _BULLET_UNI.get(sp.bullet, "•")

    def _name(run) -> None:
        run.font.name = ofont
        rpr = run._element.get_or_add_rPr()
        rfonts = rpr.find(qn("w:rFonts"))
        if rfonts is None:
            rfonts = OxmlElement("w:rFonts")
            rpr.append(rfonts)
        for attr in ("w:ascii", "w:hAnsi", "w:cs"):
            rfonts.set(qn(attr), ofont)

    def _shade_el(el, hex_fill: str, pr_getter) -> None:
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_fill)
        pr_getter().append(shd)

    def shade_cell(cell, hex_fill: str) -> None:
        _shade_el(cell, hex_fill, cell._tc.get_or_add_tcPr)

    def shade_par(par, hex_fill: str) -> None:
        _shade_el(par, hex_fill, par._p.get_or_add_pPr)

    def par_border(par, edges: dict[str, tuple[str, int]]) -> None:
        """edges: {'left': (hexcolor, eighths_pt), ...}."""
        pPr = par._p.get_or_add_pPr()
        pbdr = OxmlElement("w:pBdr")
        for side, (color, sz) in edges.items():
            e = OxmlElement(f"w:{side}")
            e.set(qn("w:val"), "single")
            e.set(qn("w:sz"), str(sz))
            e.set(qn("w:space"), "6")
            e.set(qn("w:color"), color)
            pbdr.append(e)
        pPr.append(pbdr)

    def para(text: str = "", *, size: int | None = None, color=ink, bold=False, italic=False,
             align=None, spacing: float | None = None) -> Any:
        p = doc.add_paragraph()
        if align is not None:
            p.alignment = align
        if spacing is not None:
            p.paragraph_format.line_spacing = spacing
        if text:
            r = p.add_run(text)
            r.bold, r.italic = bold, italic
            if size:
                r.font.size = Pt(size)
            r.font.color.rgb = color
            _name(r)
        return p

    def add_data_table(headers: list[str], rows: list[list[str]]) -> None:
        ts = sp.table_style
        t = doc.add_table(rows=1, cols=len(headers) or 1)
        try:
            t.style = "Table Grid" if ts in ("grid", "boxed_header") else "Light List"
        except Exception:  # noqa: BLE001
            try:
                t.style = "Table Grid"
            except Exception:  # noqa: BLE001
                pass
        t.alignment = WD_ALIGN_PARAGRAPH.LEFT
        head_fill = ts not in ("minimal", "borderless")
        for c, htext in enumerate(headers):
            cell = t.rows[0].cells[c]
            cell.text = ""
            run = cell.paragraphs[0].add_run(str(htext))
            run.bold = True
            run.font.color.rgb = RGBColor(*_WHITE) if head_fill else heading_c
            _name(run)
            if head_fill:
                shade_cell(cell, _hex(theme.primary))
        for ri, row in enumerate(rows):
            cells = t.add_row().cells
            for c, val in enumerate(row):
                cells[c].text = str(val)
                for r in cells[c].paragraphs[0].runs:
                    _name(r)
                if ts in ("zebra", "boxed_header") and ri % 2 == 1:
                    shade_cell(cells[c], _hex(theme.row_alt))

    def place_image(rel_abs: Path) -> None:
        try:
            doc.add_picture(str(rel_abs), width=Inches(5.6 if sp.image_policy == "fullbleed" else 4.4))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception as exc:  # noqa: BLE001
            logger.warning("docx image embed skipped: {}", repr(exc))

    doc = Document()
    for style_name in ("Normal", "List Bullet", "List Number", "Title", "Heading 1", "Heading 2"):
        try:
            doc.styles[style_name].font.name = ofont
        except KeyError:  # noqa: PERF203 - style absent in this template
            pass
    base = doc.styles["Normal"]
    base.paragraph_format.line_spacing = sp.line_spacing
    base.font.size = Pt(sp.body_size)

    # ---------------------------------------------------------------- cover
    cover_img: Path | None = None
    consumed_id: int | None = None
    if sp.cover in ("fullbleed_image", "magazine"):
        for s in sections:
            ap = _img_abs(s, asset_root)
            if ap is not None:
                cover_img, consumed_id = ap, id(s)
                break

    cv = sp.cover
    if cv in ("fullbleed_image", "magazine") and cover_img is not None:
        try:
            doc.add_picture(str(cover_img), width=Inches(6.5))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception:  # noqa: BLE001
            pass
    if cv == "letterhead":
        lh = para("OMNIVRA AI COMPANY OS", size=10, color=primary, bold=True)
        par_border(lh, {"bottom": (_hex(theme.accent), 6)})
    if cv == "masthead":
        mp = para(title or "Untitled", size=26, color=primary, bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
        par_border(mp, {"top": (_hex(theme.ink), 12), "bottom": (_hex(theme.ink), 12)})
        para(subtitle or sp.genre, size=9, color=muted, italic=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    else:
        title_align = WD_ALIGN_PARAGRAPH.CENTER if cv in ("centered", "fullbleed_image", "magazine") else WD_ALIGN_PARAGRAPH.LEFT
        tp = para(title or "Untitled", size=28 if cv in ("plain", "centered") else 24,
                  color=primary, bold=True, align=title_align)
        if cv == "band":
            shade_par(tp, _hex(theme.primary))
            for r in tp.runs:
                r.font.color.rgb = RGBColor(*_WHITE)
        if cv == "magazine":
            kp = para(sp.genre.upper(), size=10, color=accent, bold=True)
            tp._p.addprevious(kp._p)  # move the kicker line ABOVE the title
        if subtitle:
            para(subtitle, size=13, color=accent if cv != "band" else muted, italic=True, align=title_align)

    if cv == "centered" and sp.cover_page:
        para().add_run().add_break(WD_BREAK.PAGE)

    # body columns: a fresh section (continuous, or new-page when the cover is dedicated)
    if sp.columns == 2:
        section = doc.add_section(WD_SECTION.NEW_PAGE if sp.cover_page else WD_SECTION.CONTINUOUS)
        try:
            cols = section._sectPr.find(qn("w:cols"))
            if cols is None:
                cols = OxmlElement("w:cols")
                section._sectPr.append(cols)
            cols.set(qn("w:num"), "2")
            cols.set(qn("w:space"), "360")
        except Exception:  # noqa: BLE001
            pass
    elif sp.cover_page and cv not in ("centered",):
        para().add_run().add_break(WD_BREAK.PAGE)

    # ---------------------------------------------------------------- TOC
    if sp.toc and len(sections) > 1:
        tp = para("Contents", size=15, color=heading_c, bold=True)
        par_border(tp, {"bottom": (_hex(theme.accent), 6)})
        for i, s in enumerate(sections, start=1):
            num = _sec_number(sp.numbering, i)
            para(f"{num + '  ' if num else ''}{s.get('heading', '')}", size=11, color=ink)

    # ---------------------------------------------------------------- body
    for idx, sec in enumerate(sections, start=1):
        if sp.divider == "new_page" and idx > 1:
            para().add_run().add_break(WD_BREAK.PAGE)
        num = _sec_number(sp.numbering, idx)
        htext = f"{num}  {sec.get('heading', '')}" if num else sec.get("heading", "")
        hs = sp.heading_style
        if hs == "kicker":
            para(sp.genre.upper(), size=9, color=accent, bold=True, align=halign)
        hp = para(htext.upper() if hs == "allcaps" else htext,
                  size=14 if hs == "allcaps" else 15, color=heading_c, bold=True, align=halign)
        if hs == "rule_under":
            par_border(hp, {"bottom": (_hex(theme.accent), 6)})
        elif hs == "bottom_border":
            par_border(hp, {"bottom": (_hex(theme.accent), 18)})
        elif hs == "left_bar":
            par_border(hp, {"left": (_hex(theme.accent), 24)})
        elif hs == "boxed":
            shade_par(hp, _hex(theme.row_alt))
            par_border(hp, {"left": (_hex(theme.accent), 8), "right": (_hex(theme.accent), 8),
                            "top": (_hex(theme.accent), 8), "bottom": (_hex(theme.accent), 8)})

        body = (sec.get("body") or "").strip()
        if body:
            parts = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
            if sp.callout == "pull_quote" and parts and len(body) < 240 and not sec.get("bullets"):
                qp = para(f"“{parts[0]}”", size=sp.body_size + 5, color=primary, bold=True, italic=True,
                          align=WD_ALIGN_PARAGRAPH.CENTER)
                parts = parts[1:]
            for pi, ptext in enumerate(parts):
                bp = para(align=align_body, spacing=sp.line_spacing)
                if sp.drop_cap and idx == 1 and pi == 0 and len(ptext) > 1:
                    cap = bp.add_run(ptext[0]); cap.bold = True; cap.font.size = Pt(sp.body_size * 2)
                    cap.font.color.rgb = primary; _name(cap)
                    rest = bp.add_run(ptext[1:]); rest.font.color.rgb = ink; _name(rest)
                else:
                    r = bp.add_run(ptext); r.font.color.rgb = ink; _name(r)
                if sp.container in ("shaded", "card"):
                    shade_par(bp, _hex(theme.row_alt))
                if sp.container == "bordered":
                    par_border(bp, {"left": (_hex(theme.accent), 6)})

        for i, b in enumerate(sec.get("bullets") or [], start=1):
            text = str(b).strip()
            if not text:
                continue
            mark = f"{i}." if sp.bullet == "numbered" else glyph
            bp = doc.add_paragraph()
            bp.paragraph_format.left_indent = Inches(0.25)
            if sp.bullet != "none":
                mr = bp.add_run(f"{mark}  "); mr.font.color.rgb = accent; mr.bold = True; _name(mr)
            tr = bp.add_run(text); tr.font.color.rgb = ink; _name(tr)

        # callouts that wrap a dedicated note (warning/speech/quote_bar/tint) apply to a short body line
        if sp.callout in ("warning", "speech") and not body and sec.get("bullets"):
            pass  # already rendered as bullets

        table = _norm_table(sec.get("table"))
        if table:
            add_data_table(*table)
        chart = _norm_chart(sec.get("chart"))
        if chart:
            caption, (c_headers, c_rows) = _chart_to_table(chart)
            para(caption, size=11, color=heading_c, bold=True, italic=True)
            add_data_table(c_headers, c_rows)

        ap = _img_abs(sec, asset_root)
        if ap is not None and id(sec) != consumed_id and sp.image_policy != "none":
            place_image(ap)

        if sp.divider == "thick":
            dp = para(); par_border(dp, {"bottom": (_hex(theme.accent), 18)})
        elif sp.divider == "shaded":
            dp = para(); shade_par(dp, _hex(theme.row_alt))
        elif sp.divider == "ornament":
            para("·  ·  ·", size=12, color=accent, align=WD_ALIGN_PARAGRAPH.CENTER)
        elif sp.divider == "hairline":
            dp = para(); par_border(dp, {"bottom": (_hex(theme.row_alt), 4)})

    # ---------------------------------------------------------------- CTA
    if sp.cta:
        cp = para(subtitle or "Ready to take the next step?", size=15, color=RGBColor(*_WHITE),
                  bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
        shade_par(cp, _hex(theme.primary))

    # ---------------------------------------------- page frame: borders + header/footer
    for section in doc.sections:
        if sp.page_border in ("full", "top_bottom"):
            try:
                sectPr = section._sectPr
                pgB = sectPr.find(qn("w:pgBorders"))
                if pgB is None:
                    pgB = OxmlElement("w:pgBorders")
                    pgB.set(qn("w:offsetFrom"), "page")
                    # w:pgBorders must precede w:cols/w:docGrid in the CT_SectPr sequence; a plain
                    # append() lands it last -> schema-invalid -> Word repairs the file & drops the frame.
                    sectPr.insert_element_before(
                        pgB, "w:lnNumType", "w:pgNumType", "w:cols", "w:formProt", "w:vAlign",
                        "w:noEndnote", "w:titlePg", "w:textDirection", "w:bidi", "w:rtlGutter",
                        "w:docGrid", "w:printerSettings", "w:sectPrChange",
                    )
                sides = ("top", "left", "bottom", "right") if sp.page_border == "full" else ("top", "bottom")
                for side in sides:
                    e = OxmlElement(f"w:{side}")
                    e.set(qn("w:val"), "single")
                    e.set(qn("w:sz"), "8")
                    e.set(qn("w:space"), "24")
                    e.set(qn("w:color"), _hex(theme.accent))
                    pgB.append(e)
            except Exception:  # noqa: BLE001
                pass
        if sp.header in ("running_head", "masthead", "letterhead"):
            try:
                hp = section.header.paragraphs[0]
                hp.text = ""
                label = (title or "")[:80] if sp.header != "letterhead" else sp.genre.upper()
                r = hp.add_run(label)
                r.font.size = Pt(8); r.italic = sp.header == "running_head"; r.font.color.rgb = muted
                _name(r)
            except Exception:  # noqa: BLE001
                pass
        if sp.footer in ("page_num", "page_x_of_y", "brand"):
            try:
                fp = section.footer.paragraphs[0]
                fp.text = ""
                fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if sp.footer == "brand":
                    r = fp.add_run("Omnivra AI Company OS"); r.font.size = Pt(8); r.font.color.rgb = muted; _name(r)
                else:
                    _docx_page_field(fp, with_total=sp.footer == "page_x_of_y", color=muted, name=ofont)
            except Exception:  # noqa: BLE001
                pass

    doc.save(str(out))


def _docx_page_field(paragraph, *, with_total: bool, color, name: str) -> None:
    """Insert a PAGE (and optionally 'of NUMPAGES') field into a footer paragraph."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt

    def field(code: str):
        run = paragraph.add_run()
        run.font.size = Pt(8)
        run.font.color.rgb = color
        run.font.name = name
        begin = OxmlElement("w:fldChar"); begin.set(qn("w:fldCharType"), "begin")
        instr = OxmlElement("w:instrText"); instr.set(qn("xml:space"), "preserve"); instr.text = f" {code} "
        end = OxmlElement("w:fldChar"); end.set(qn("w:fldCharType"), "end")
        run._r.append(begin); run._r.append(instr); run._r.append(end)

    if with_total:
        pre = paragraph.add_run("Page "); pre.font.size = Pt(8); pre.font.color.rgb = color; pre.font.name = name
        field("PAGE")
        mid = paragraph.add_run(" of "); mid.font.size = Pt(8); mid.font.color.rgb = color; mid.font.name = name
        field("NUMPAGES")
    else:
        field("PAGE")


# ===========================================================================
# PDF
# ===========================================================================
def _img_abs(sec: dict[str, Any], asset_root: Path | None) -> Path | None:
    """Resolve a section's generated image to an existing absolute path (or None)."""
    im = sec.get("image")
    if not isinstance(im, dict) or asset_root is None:
        return None
    rel = im.get("path")
    if not rel:
        return None
    try:
        ap = asset_root / str(rel)
        return ap if ap.exists() and ap.stat().st_size > 0 else None
    except Exception:  # noqa: BLE001
        return None


def _render_pdf(
    title: str, subtitle: str, sections: list[dict[str, Any]], out: Path, theme: Theme,
    *, font: str = "sans", spec: "StructureSpec | None" = None, asset_root: Path | None = None,
) -> None:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    sp = spec or _DEFAULT_STRUCT
    fam = pdf_font(font)
    m = _MARGIN_MM.get(sp.margin, 16.0)
    top_m = 22.0 if sp.header != "none" else 16.0
    EPW = 210 - 2 * m
    PAGE_BOTTOM = 297 - 16
    accent_rgb, primary_rgb, heading_rgb, ink_rgb = theme.accent, theme.primary, theme.heading, theme.ink
    muted_rgb = (130, 130, 138)
    bullet_glyph = _BULLET_PDF.get(sp.bullet, "\xb7")

    class _Doc(FPDF):
        cover_on_first = False
        suppress = False

        def header(self) -> None:  # runs on every add_page (incl. content pages)
            if self.suppress:
                return
            if sp.page_border == "full":
                self.set_draw_color(*accent_rgb)
                self.set_line_width(0.4)
                self.rect(8, 8, 210 - 16, 297 - 16)
            elif sp.page_border == "top_bottom":
                self.set_draw_color(*accent_rgb)
                self.set_line_width(0.5)
                self.line(m, 11, 210 - m, 11)
                self.line(m, 286, 210 - m, 286)
            if not (self.cover_on_first and self.page_no() == 1) and sp.header != "none":
                self.set_y(11)
                self.set_text_color(*muted_rgb)
                if sp.header == "masthead":
                    self.set_font(fam, "B", 9)
                    self.set_x(m)
                    self.cell(EPW, 5, _latin1((title or "")[:80]).upper(), align="C")
                elif sp.header == "letterhead":
                    self.set_font(fam, "B", 9)
                    self.set_x(m)
                    self.cell(EPW, 5, _latin1((sp.genre or "Document")).upper(), align="L")
                else:  # running_head
                    self.set_font(fam, "I", 8)
                    self.set_x(m)
                    self.cell(EPW, 5, _latin1((title or "")[:70]), align="L")
                self.set_draw_color(*accent_rgb)
                self.set_line_width(0.3)
                self.line(m, 17, 210 - m, 17)
            self.set_y(top_m)

        def footer(self) -> None:
            if self.suppress or sp.footer == "none":
                return
            self.set_y(-13)
            self.set_text_color(*muted_rgb)
            self.set_font(fam, "", 8)
            if sp.footer == "page_x_of_y":
                txt = f"Page {self.page_no()} of {{nb}}"
            elif sp.footer == "brand":
                txt = "Omnivra AI Company OS"
            else:  # page_num
                txt = str(self.page_no())
            self.cell(0, 6, _latin1(txt), align="C")

    pdf = _Doc()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(m, top_m, m)
    try:
        pdf.set_title(_latin1(title or "Document"))
    except Exception:  # noqa: BLE001
        pass

    def txt(s: str) -> str:
        return _latin1(s or "")

    def rule(y: float, color: RGB, w: float = 0.5, x0: float | None = None, x1: float | None = None) -> None:
        pdf.set_draw_color(*color)
        pdf.set_line_width(w)
        pdf.line(x0 if x0 is not None else m, y, x1 if x1 is not None else 210 - m, y)

    # ---------------------------------------------------------------- cover
    cover_img: Path | None = None
    consumed_img_id: int | None = None
    if sp.cover in ("fullbleed_image", "magazine"):
        for i, s in enumerate(sections):
            ap = _img_abs(s, asset_root)
            if ap is not None:
                cover_img, consumed_img_id = ap, id(s)
                break

    dedicated = sp.cover_page or sp.cover in ("fullbleed_image", "magazine", "centered")

    def draw_cover() -> None:
        cv = sp.cover
        if cv == "fullbleed_image":
            if cover_img is not None:
                try:
                    pdf.image(str(cover_img), x=0, y=0, w=210, h=297)
                except Exception:  # noqa: BLE001
                    pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 297, "F")
            else:
                pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 297, "F")
            pdf.set_fill_color(*ink_rgb); pdf.rect(0, 196, 210, 101, "F")  # overlay band
            pdf.set_xy(m, 214); pdf.set_text_color(*_WHITE); pdf.set_font(fam, "B", 34)
            pdf.multi_cell(EPW, 13, txt(title or "Untitled"), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_font(fam, "", 13)
                pdf.multi_cell(EPW, 7, txt(subtitle), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        elif cv == "magazine":
            if cover_img is not None:
                try:
                    pdf.image(str(cover_img), x=0, y=0, w=210, h=168)
                except Exception:  # noqa: BLE001
                    pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 168, "F")
            else:
                pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 168, "F")
            pdf.set_xy(m, 182); pdf.set_text_color(*accent_rgb); pdf.set_font(fam, "B", 11)
            pdf.cell(EPW, 6, txt((sp.genre or "Feature").upper()))
            pdf.set_xy(m, 192); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 36)
            pdf.multi_cell(EPW, 14, txt(title or "Untitled"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "I", 13)
                pdf.multi_cell(EPW, 7, txt(subtitle), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        elif cv == "centered":
            pdf.set_y(110); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 30)
            pdf.multi_cell(0, 13, txt(title or "Untitled"), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            ry = pdf.get_y() + 3; rule(ry, accent_rgb, 0.8, 80, 130); pdf.ln(7)
            if subtitle:
                pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "I", 13)
                pdf.multi_cell(0, 7, txt(subtitle), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_y(272); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "", 9)
            pdf.cell(0, 6, txt(f"{sp.genre} \xb7 Omnivra AI Company OS"), align="C")
        elif cv == "letterhead":
            pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 4, "F")
            pdf.set_xy(m, 14); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 11)
            pdf.cell(EPW, 6, txt("OMNIVRA AI COMPANY OS"), align="L")
            pdf.set_xy(m, 14); pdf.set_font(fam, "", 9); pdf.set_text_color(*muted_rgb)
            pdf.cell(EPW, 6, txt(sp.genre.upper()), align="R")
            rule(24, muted_rgb, 0.3)
            yy = 150 if dedicated else 32
            pdf.set_xy(m, yy); pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "B", 22)
            pdf.multi_cell(EPW, 10, txt(title or "Untitled"), align="C" if dedicated else "L",
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "I", 12)
                pdf.multi_cell(EPW, 6, txt(subtitle), align="C" if dedicated else "L",
                               new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        elif cv == "band":
            pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 210, 38, "F")
            rule(38, accent_rgb, 1.4, 0, 210)
            pdf.set_xy(m, 11); pdf.set_text_color(*_WHITE); pdf.set_font(fam, "B", 22)
            pdf.multi_cell(EPW, 9, txt(title or "Untitled"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_font(fam, "", 11)
                pdf.multi_cell(EPW, 6, txt(subtitle), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_y(46)
        elif cv == "masthead":
            rule(12, ink_rgb, 1.0, 0, 210)
            pdf.set_xy(0, 15); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 30)
            pdf.multi_cell(210, 12, txt(title or "Untitled"), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            yy = pdf.get_y() + 1; rule(yy, ink_rgb, 1.0, 0, 210)
            pdf.set_xy(m, yy + 2); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "I", 9)
            pdf.cell(EPW, 5, txt(subtitle or sp.genre), align="C")
            pdf.set_y(yy + 11)
        elif cv == "hero_number":
            pdf.set_fill_color(*primary_rgb); pdf.rect(0, 0, 14, 60, "F")
            pdf.set_xy(m, 14); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 26)
            pdf.multi_cell(EPW, 11, txt(title or "Untitled"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "I", 12)
                pdf.multi_cell(EPW, 6, txt(subtitle), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(2); rule(pdf.get_y(), accent_rgb, 1.2); pdf.ln(4)
        elif cv == "inline":
            pdf.set_xy(m, top_m); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 19)
            pdf.multi_cell(EPW, 9, txt(title or "Untitled"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "I", 11)
                pdf.multi_cell(EPW, 6, txt(subtitle), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(1); rule(pdf.get_y(), accent_rgb, 0.6); pdf.ln(3)
        else:  # plain
            pdf.set_xy(m, top_m); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "B", 30)
            pdf.multi_cell(EPW, 12, txt(title or "Untitled"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            if subtitle:
                pdf.set_x(m); pdf.set_text_color(*muted_rgb); pdf.set_font(fam, "", 13)
                pdf.multi_cell(EPW, 7, txt(subtitle), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(3)

    pdf.cover_on_first = not dedicated and sp.cover != "inline"
    pdf.suppress = dedicated
    pdf.add_page()
    draw_cover()
    if dedicated:
        pdf.suppress = False
        pdf.add_page()

    # ---------------------------------------------------------------- TOC
    if sp.toc and len(sections) > 1:
        pdf.set_text_color(*heading_rgb); pdf.set_font(fam, "B", 16)
        pdf.set_x(m); pdf.multi_cell(EPW, 9, txt("Contents"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        rule(pdf.get_y() + 1, accent_rgb, 0.5); pdf.ln(4)
        for i, s in enumerate(sections, start=1):
            label = s.get("heading", "")
            num = _sec_number(sp.numbering, i)
            pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "", 11)
            pdf.set_x(m + 2)
            pdf.multi_cell(EPW - 2, 7, txt(f"{num + '  ' if num else ''}{label}"),
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)
        if not dedicated:
            pdf.add_page()

    # ------------------------------------------------------ section helpers
    def heading(text: str, idx: int) -> None:
        num = _sec_number(sp.numbering, idx)
        label = f"{num}  {text}" if num else text
        align = "C" if sp.heading_align == "center" else "L"
        hs = sp.heading_style
        if pdf.get_y() > PAGE_BOTTOM - 30 and sp.divider != "new_page":
            pdf.add_page()
        pdf.set_text_color(*heading_rgb)
        if hs == "kicker":
            pdf.set_font(fam, "B", 9); pdf.set_text_color(*accent_rgb); pdf.set_x(m)
            pdf.cell(EPW, 5, txt((sp.genre or "Section").upper()), align=align)
            pdf.ln(5); pdf.set_text_color(*heading_rgb)
        if hs == "boxed":
            y0 = pdf.get_y()
            pdf.set_fill_color(*theme.row_alt); pdf.set_draw_color(*accent_rgb); pdf.set_line_width(0.4)
            pdf.set_xy(m, y0); pdf.set_font(fam, "B", 14)
            pdf.multi_cell(EPW, 9, txt(label), border=1, fill=True, align=align,
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(2); return
        if hs == "left_bar":
            y0 = pdf.get_y()
            pdf.set_fill_color(*accent_rgb); pdf.rect(m, y0 + 0.5, 1.6, 7, "F")
            pdf.set_xy(m + 4, y0); pdf.set_font(fam, "B", 15)
            pdf.multi_cell(EPW - 4, 8, txt(label), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(2); return
        size = 14 if hs == "allcaps" else 15
        pdf.set_font(fam, "B", size); pdf.set_x(m)
        shown = txt(label).upper() if hs == "allcaps" else txt(label)
        pdf.multi_cell(EPW, 8, shown, align=align, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        if hs == "rule_under":
            rule(pdf.get_y() + 0.5, accent_rgb, 0.6); pdf.ln(2.5)
        elif hs == "bottom_border":
            rule(pdf.get_y() + 0.5, accent_rgb, 1.6, m, m + 26); pdf.ln(3)
        else:
            pdf.ln(1.5)

    def paragraphs(body: str, drop: bool) -> None:
        parts = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
        align = "J" if sp.justify else "L"
        lh = sp.body_size * 0.42 * sp.line_spacing
        if sp.columns == 2 and parts:
            if _pdf_two_columns(pdf, parts, m, EPW, fam, sp.body_size, lh, ink_rgb, align, top_m, PAGE_BOTTOM):
                pdf.ln(1.5); return
        for pi, para in enumerate(parts):
            pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "", sp.body_size)
            if drop and pi == 0 and len(para) > 1:
                cap, rest = para[0], para[1:].lstrip()
                y0 = pdf.get_y()
                pdf.set_font(fam, "B", sp.body_size * 2 + 8); pdf.set_text_color(*primary_rgb)
                cw = pdf.get_string_width(cap) + 2
                pdf.set_xy(m, y0); pdf.cell(cw, lh * 2, txt(cap))
                pdf.set_xy(m + cw, y0); pdf.set_font(fam, "", sp.body_size); pdf.set_text_color(*ink_rgb)
                pdf.multi_cell(EPW - cw, lh, txt(rest), align=align, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            else:
                pdf.set_x(m)
                pdf.multi_cell(EPW, lh, txt(para), align=align, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(1)

    def bullets(items: list[str]) -> None:
        lh = sp.body_size * 0.42 * sp.line_spacing
        for i, b in enumerate(items, start=1):
            text = str(b).strip()
            if not text:
                continue
            mark = f"{i}." if sp.bullet == "numbered" else bullet_glyph
            if sp.bullet == "none":
                mark = ""
            pdf.set_text_color(*accent_rgb); pdf.set_font(fam, "B", sp.body_size)
            pdf.set_x(m + 2)
            mw = 7 if sp.bullet == "numbered" else 5
            if mark:
                pdf.cell(mw, lh, txt(mark))
            pdf.set_text_color(*ink_rgb); pdf.set_font(fam, "", sp.body_size)
            pdf.multi_cell(EPW - mw - 2, lh, txt(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)

    def callout(text: str) -> None:
        style = sp.callout
        if style in ("none", ""):
            paragraphs(text, False); return
        if style == "pull_quote":
            pdf.ln(2); pdf.set_text_color(*primary_rgb); pdf.set_font(fam, "BI", sp.body_size + 5)
            pdf.set_x(m + 8)
            pdf.multi_cell(EPW - 16, sp.body_size * 0.6, txt(f"“{text}”"), align="C",
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.ln(2); return
        if style == "quote_bar":
            y0 = pdf.get_y(); pdf.set_font(fam, "I", sp.body_size)
            pdf.set_xy(m + 5, y0); pdf.set_text_color(*ink_rgb)
            pdf.multi_cell(EPW - 7, sp.body_size * 0.5, txt(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_fill_color(*accent_rgb); pdf.rect(m, y0, 1.4, pdf.get_y() - y0, "F"); pdf.ln(2); return
        # tint / warning / speech -> a filled box
        y0 = pdf.get_y()
        fill = theme.row_alt if style != "warning" else (255, 244, 230)
        edge = accent_rgb if style != "warning" else (217, 119, 6)
        label = {"warning": "NOTE", "speech": ""}.get(style, "")
        pdf.set_fill_color(*fill); pdf.set_draw_color(*edge); pdf.set_line_width(0.4)
        pdf.set_xy(m, y0); pdf.set_font(fam, "" if not label else "B", sp.body_size)
        pdf.set_text_color(*ink_rgb)
        content = (label + ": " if label else "") + text
        pdf.multi_cell(EPW, sp.body_size * 0.5, txt(content), border=1, fill=True,
                       new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

    def place_image(ap: Path) -> None:
        try:
            if pdf.get_y() > PAGE_BOTTOM - 60:
                pdf.add_page()
            w = min(EPW, 120 if sp.image_policy != "fullbleed" else EPW)
            y0 = pdf.get_y() + 1
            # fpdf2 does NOT auto-page-break a positioned image; cap its box to the room left and let
            # keep_aspect_ratio scale it so a tall/square image can never spill off the page bottom.
            h_box = min(w, PAGE_BOTTOM - y0)
            x = m + (EPW - w) / 2
            pdf.image(str(ap), x=x, y=y0, w=w, h=h_box, keep_aspect_ratio=True)
            pdf.set_y(y0 + h_box)
            pdf.ln(2)
        except Exception as exc:  # noqa: BLE001 - never fail the file over an image
            logger.warning("pdf image embed skipped: {}", repr(exc))

    def divider() -> None:
        d = sp.divider
        if d == "none":
            pdf.ln(3); return
        if d == "new_page":
            pdf.add_page(); return
        if d == "thick":
            pdf.ln(2); rule(pdf.get_y(), accent_rgb, 1.6); pdf.ln(4); return
        if d == "shaded":
            y0 = pdf.get_y() + 1; pdf.set_fill_color(*theme.row_alt); pdf.rect(m, y0, EPW, 2, "F"); pdf.ln(5); return
        if d == "ornament":
            pdf.ln(1); pdf.set_text_color(*accent_rgb); pdf.set_font(fam, "B", 12)
            pdf.cell(0, 6, txt("\xb7  \xb7  \xb7"), align="C", new_x=XPos.LMARGIN, new_y=YPos.NEXT); pdf.ln(2); return
        pdf.ln(2); rule(pdf.get_y(), muted_rgb, 0.2); pdf.ln(3)  # hairline

    # ---------------------------------------------------------------- body
    container_box = sp.container in ("card", "shaded", "bordered")
    for idx, sec in enumerate(sections, start=1):
        box_y0 = pdf.get_y()
        heading(sec.get("heading", ""), idx)
        body = (sec.get("body") or "").strip()
        if body:
            if sp.callout == "pull_quote" and len(body) < 240 and not sec.get("bullets"):
                callout(body)
            else:
                paragraphs(body, sp.drop_cap and idx == 1)
        for_items = [str(b).strip() for b in (sec.get("bullets") or []) if str(b).strip()]
        if for_items:
            bullets(for_items)
        table = _norm_table(sec.get("table"))
        if table:
            _pdf_table(pdf, table, theme, EPW, fam, sp.table_style, m)
        chart = _norm_chart(sec.get("chart"))
        if chart:
            caption, c_table = _chart_to_table(chart)
            pdf.set_text_color(*heading_rgb); pdf.set_font(fam, "BI", 11); pdf.set_x(m)
            pdf.multi_cell(EPW, 6, txt(caption), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            _pdf_table(pdf, c_table, theme, EPW, fam, sp.table_style, m)
        ap = _img_abs(sec, asset_root)
        if ap is not None and id(sec) != consumed_img_id and sp.image_policy not in ("none",):
            place_image(ap)
        if container_box:
            try:  # frame the whole section block
                edge = theme.row_alt if sp.container == "shaded" else accent_rgb
                pdf.set_draw_color(*edge); pdf.set_line_width(0.3)
                pdf.rect(m - 2, box_y0 - 1, EPW + 4, pdf.get_y() - box_y0 + 2)
            except Exception:  # noqa: BLE001
                pass
        if idx < len(sections):  # no divider after the last section (avoids a trailing rule / blank page)
            divider()

    # ---------------------------------------------------------------- CTA
    if sp.cta:
        if pdf.get_y() > PAGE_BOTTOM - 40:
            pdf.add_page()
        y0 = pdf.get_y() + 2
        pdf.set_fill_color(*primary_rgb); pdf.rect(m, y0, EPW, 26, "F")
        pdf.set_xy(m, y0 + 6); pdf.set_text_color(*_WHITE); pdf.set_font(fam, "B", 15)
        pdf.multi_cell(EPW, 8, txt(subtitle or "Ready to take the next step?"), align="C",
                       new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.output(str(out))


def _pdf_two_columns(pdf, parts, m, epw, fam, size, lh, ink, align, top_m, page_bottom) -> bool:
    """Flow body paragraphs into TWO balanced columns (magazine/journal look). Manual placement with
    overflow handling. Returns True on success; False (caller falls back to single column) on any error."""
    from fpdf.enums import XPos, YPos

    try:
        gutter = 7.0
        colw = (epw - gutter) / 2
        xs = [m, m + colw + gutter]
        top = pdf.get_y()
        ys = [top, top]
        col = 0
        pdf.set_auto_page_break(False)
        for para in parts:
            pdf.set_font(fam, "", size)
            lines = pdf.multi_cell(colw, lh, _latin1(para), align=align, split_only=True)
            need = max(lh, len(lines) * lh) + 2
            if ys[col] + need > page_bottom:
                if col == 0:
                    col = 1
                else:
                    pdf.set_auto_page_break(True, margin=18)
                    pdf.add_page()
                    pdf.set_auto_page_break(False)
                    top = pdf.get_y(); ys = [top, top]; col = 0
            pdf.set_xy(xs[col], ys[col])
            pdf.set_text_color(*ink)
            # YPos.NEXT (not TOP) so get_y() lands at the paragraph's BOTTOM — TOP leaves y unchanged,
            # which would advance the column cursor ~0mm and stack every paragraph on the previous one.
            pdf.multi_cell(colw, lh, _latin1(para), align=align, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            ys[col] = pdf.get_y() + 2
        pdf.set_auto_page_break(True, margin=18)
        pdf.set_y(max(ys))
        return True
    except Exception as exc:  # noqa: BLE001 - fall back to single-column prose
        logger.warning("pdf 2-column flow fell back to single column: {}", repr(exc))
        try:
            pdf.set_auto_page_break(True, margin=18)
        except Exception:  # noqa: BLE001
            pass
        return False


def _sec_number(scheme: str, idx: int) -> str:
    """Section number string for a numbering scheme (decimal/legal/section_word/none)."""
    if scheme == "decimal":
        return f"{idx}."
    if scheme == "legal":
        return f"{idx}.0"
    if scheme == "section_word":
        return f"Section {idx}"
    return ""


def _pdf_table(pdf, table: tuple[list[str], list[list[str]]], theme: Theme, epw: float,
               fam: str = "Helvetica", table_style: str = "zebra", m: float = 15.0) -> None:
    """Draw a styled table honoring ``table_style`` (grid|zebra|minimal|boxed_header|borderless).
    Falls back to plain lines on error so a section is never lost."""
    from fpdf.enums import XPos, YPos

    headers, rows = table
    ncols = len(headers) or 1
    col_w = epw / ncols
    line_h = 7
    border_h = 1 if table_style in ("grid", "boxed_header") else 0
    zebra = table_style in ("zebra", "boxed_header")
    try:
        pdf.set_font(fam, "B", 10)
        if table_style == "minimal" or table_style == "borderless":
            pdf.set_text_color(*theme.heading)
        else:
            pdf.set_fill_color(*theme.primary); pdf.set_text_color(*_WHITE)
        pdf.set_draw_color(200, 200, 205)
        pdf.set_x(m)
        head_fill = table_style not in ("minimal", "borderless")
        for htext in headers:
            pdf.cell(col_w, line_h, _latin1(str(htext))[:40], border=border_h, fill=head_fill, align="L")
        pdf.ln(line_h)
        if table_style == "minimal":  # rule under header only
            pdf.set_draw_color(*theme.accent); pdf.set_line_width(0.4)
            pdf.line(m, pdf.get_y(), m + epw, pdf.get_y())
        pdf.set_font(fam, "", 10); pdf.set_text_color(*theme.ink)
        for ri, row in enumerate(rows):
            if pdf.get_y() + line_h > 282:
                pdf.add_page()
            fill = zebra and ri % 2 == 1
            if fill:
                pdf.set_fill_color(*theme.row_alt)
            pdf.set_x(m)
            for val in row:
                pdf.cell(col_w, line_h, _latin1(str(val))[:40], border=border_h, fill=fill, align="L")
            pdf.ln(line_h)
            if table_style == "minimal":
                pdf.set_draw_color(235, 235, 238); pdf.set_line_width(0.2)
                pdf.line(m, pdf.get_y(), m + epw, pdf.get_y())
        pdf.ln(2)
    except Exception as exc:  # noqa: BLE001 - degrade to plain text, never lose the section
        logger.warning("pdf table render fell back to text: {}", repr(exc))
        pdf.set_text_color(*theme.ink)
        pdf.set_font(fam, "", 10)
        for row in [headers, *rows]:
            pdf.set_x(m)
            pdf.multi_cell(0, 6, _latin1(" | ".join(str(c) for c in row)), new_x=XPos.LMARGIN, new_y=YPos.NEXT)


_RENDERERS = {"pptx": _render_pptx, "docx": _render_docx, "pdf": _render_pdf}


def render_document(
    title: str,
    sections: list[dict[str, Any]],
    output_path: Path,
    fmt: str,
    *,
    subtitle: str = "",
    theme: str = _DEFAULT_THEME,
    style: str = "professional",
    font: str = "sans",
    asset_root: Path | None = None,
) -> dict[str, Any]:
    """Render ``title`` + ``sections`` to ``output_path`` in ``fmt``.

    Three INDEPENDENT axes: ``style`` selects the document GENRE (its visual structure — cover,
    columns, numbering, page frame, headings, dividers, callouts); ``theme`` is the user's color
    palette (unknown/'auto' -> the default, never derived from the genre); ``font`` is the user's
    typeface family (serif|sans|mono). Section ``image`` paths (resolved under ``asset_root``) are
    embedded per the genre's image policy. Returns {ok, stub, path?, note}. Never raises.
    """
    if os.environ.get("OMNIVRA_DISABLE_RENDER"):
        return {"ok": True, "stub": True, "note": "render engine disabled (OMNIVRA_DISABLE_RENDER); markdown saved as the draft."}
    renderer = _RENDERERS.get(fmt)
    if renderer is None:
        return {"ok": False, "stub": False, "note": f"Unsupported format {fmt!r}."}
    try:
        __import__(_FMT_MODULE[fmt])
    except Exception as exc:  # noqa: BLE001 - lib not installed -> stub
        return {"ok": True, "stub": True, "note": f"{fmt} engine not installed ({type(exc).__name__}); markdown saved. pip install -r requirements-docs.txt"}
    try:
        spec = structure_spec(style)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        renderer(title, subtitle, sections, output_path, resolve_theme(theme),
                 font=font, spec=spec, asset_root=asset_root)
        return {"ok": True, "stub": False, "path": str(output_path), "note": f"Rendered {fmt}."}
    except Exception as exc:  # noqa: BLE001 - report a render failure without raising
        logger.error("Document render failed ({}): {}", fmt, repr(exc))
        return {"ok": False, "stub": False, "note": f"{fmt} render failed: {type(exc).__name__}"}
