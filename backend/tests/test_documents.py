"""Document Studio tests (cp-0025): generate -> approve/reject, formats, isolation.

Everything runs offline: with no provider keys the documentation agent (Gemma) stubs,
so content falls back to the deterministic builder; with OMNIVRA_DISABLE_RENDER=1 (set
by conftest) the render engine is forced off, so every generate yields a markdown
deliverable (stub=True). A skipif test exercises a real PPTX/DOCX/PDF when the optional
libs are installed. Project-scoped via the X-Project-Id header.
"""
from __future__ import annotations

import importlib.util

import pytest
from fastapi.testclient import TestClient

from app.services.doc_render import _FMT_MODULE, render_document


def _engine_installed(fmt: str) -> bool:
    """True if the render lib for ``fmt`` is importable (ignores the disable flag)."""
    return importlib.util.find_spec(_FMT_MODULE[fmt]) is not None


def _generate(client: TestClient, prompt: str, fmt: str | None = None) -> dict:
    """Fire-and-poll: POST returns a 'generating' draft; return the TERMINAL draft.

    The TestClient runs the BackgroundTask before the next request, so the GET sees the
    finished (awaiting_approval) draft."""
    payload: dict = {"prompt": prompt}
    if fmt:
        payload["format"] = fmt
    res = client.post("/api/documents/generate", json=payload)
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["status"] == "generating", body
    return client.get(f"/api/documents/{body['id']}").json()


def test_generate_document_offline_stub(client: TestClient) -> None:
    body = _generate(client, "A roadmap for our AI company OS")
    assert body["status"] == "awaiting_approval"
    assert body["format"] == "pdf"  # default format
    assert body["title"], "a document must have a title"
    assert body["sections"], "a document must have sections"
    # Render engine is disabled in tests -> markdown deliverable.
    assert body["stub"] is True
    assert body["filePath"].endswith("document.md")
    assert any(a.endswith("content.json") for a in body["artifacts"])
    assert any(a.endswith("document.md") for a in body["artifacts"])
    # The stub explanation lives in renderNote (NOT the decision `note`), which is unset at draft time.
    assert body["renderNote"], "a stub draft must explain why it is a markdown deliverable"
    assert body["note"] is None


@pytest.mark.parametrize("fmt", ["pptx", "docx", "pdf"])
def test_generate_each_format(client: TestClient, fmt: str) -> None:
    body = _generate(client, f"Test {fmt}", fmt)
    assert body["format"] == fmt
    assert body["status"] == "awaiting_approval"


def test_generate_with_writing_style(client: TestClient) -> None:
    """A chosen writing tone rides through to the draft (placeholder + terminal)."""
    res = client.post("/api/documents/generate", json={"prompt": "a guide", "style": "legal"})
    assert res.status_code == 200
    body = res.json()
    assert body["style"] == "legal"  # the 'generating' placeholder already carries it
    done = client.get(f"/api/documents/{body['id']}").json()
    assert done["style"] == "legal"


def test_default_style_is_professional(client: TestClient) -> None:
    body = client.post("/api/documents/generate", json={"prompt": "x"}).json()
    assert body["style"] == "professional"


def test_invalid_style_rejected(client: TestClient) -> None:
    assert client.post("/api/documents/generate", json={"prompt": "x", "style": "jazzy"}).status_code == 422


def test_resolve_style_and_guide_coverage() -> None:
    from app.services.documents import _STYLE_GUIDE, _resolve_style

    assert _resolve_style("LEGAL") == "legal" and _resolve_style("nope") == "professional"
    assert len(_STYLE_GUIDE) == 20  # every offered style has a distinct instruction


def test_unknown_format_defaults_to_pdf(client: TestClient) -> None:
    # An out-of-enum format is rejected by the request schema (422), never silently mis-rendered.
    res = client.post("/api/documents/generate", json={"prompt": "x", "format": "xlsx"})
    assert res.status_code == 422


def test_approve_document(client: TestClient) -> None:
    draft = _generate(client, "Approve me")
    res = client.post(f"/api/documents/{draft['id']}/decision", json={"action": "approve"})
    assert res.status_code == 200
    body = res.json()
    assert body["status"] == "approved"
    # Approving with no note must NOT wipe the stub explanation (render_note is separate from note).
    assert body["renderNote"] == draft["renderNote"]


def test_reject_document(client: TestClient) -> None:
    draft = _generate(client, "Reject me")
    res = client.post(f"/api/documents/{draft['id']}/decision", json={"action": "reject", "note": "off-topic"})
    assert res.status_code == 200
    body = res.json()
    assert body["status"] == "rejected"
    assert body["note"] == "off-topic"
    # The rejection reason lands in `note` only; the stub explanation in render_note is preserved.
    assert body["renderNote"] == draft["renderNote"]


def test_get_and_list_document(client: TestClient) -> None:
    draft = _generate(client, "Find me later")
    got = client.get(f"/api/documents/{draft['id']}")
    assert got.status_code == 200 and got.json()["id"] == draft["id"]
    listed = client.get("/api/documents").json()
    assert any(d["id"] == draft["id"] for d in listed)


def test_documents_project_isolated(client: TestClient) -> None:
    a = client.post("/api/projects", json={"name": "Docs A"}).json()["id"]
    b = client.post("/api/projects", json={"name": "Docs B"}).json()["id"]

    client.post("/api/documents/generate", json={"prompt": "Project A doc"}, headers={"X-Project-Id": a})

    list_a = client.get("/api/documents", headers={"X-Project-Id": a}).json()
    list_b = client.get("/api/documents", headers={"X-Project-Id": b}).json()
    assert list_a and all(d["projectId"] == a for d in list_a)
    assert list_b == [], "project B must not see project A's documents (isolation)"


def test_unknown_document_404(client: TestClient) -> None:
    assert client.get("/api/documents/nope-xyz").status_code == 404
    assert client.post("/api/documents/nope-xyz/decision", json={"action": "approve"}).status_code == 404


def test_project_delete_cascades_documents(client: TestClient) -> None:
    from app.workspace_fs.paths import project_root

    pid = client.post("/api/projects", json={"name": "Docs Cascade"}).json()["id"]
    hdr = {"X-Project-Id": pid}
    client.post("/api/documents/generate", json={"prompt": "doomed doc"}, headers=hdr)
    assert client.get("/api/documents", headers=hdr).json(), "doc exists before delete"
    assert (project_root(pid) / ".state" / "documents").exists()

    assert client.delete(f"/api/projects/{pid}").status_code == 200
    assert not project_root(pid).exists(), "the project workspace (incl. documents) is purged"


def test_render_document_disabled_returns_stub(tmp_path) -> None:
    # With OMNIVRA_DISABLE_RENDER=1 (conftest) the engine is off regardless of installed libs.
    out = tmp_path / "doc.pdf"
    result = render_document("T", [{"heading": "H", "body": "B"}], out, "pdf")
    assert result["ok"] and result["stub"] and not out.exists()


# --- Content parsing: full + truncated recovery (the 'half-completed document' fix) ---
def test_parse_full_json_document() -> None:
    from app.services.documents import DocumentService

    text = '{"title":"T","subtitle":"S","theme":"violet","sections":[{"heading":"H","body":"B","bullets":["x","y"]}]}'
    parsed = DocumentService._parse(text)
    assert parsed is not None
    title, subtitle, theme, sections = parsed
    assert (title, subtitle, theme) == ("T", "S", "violet")
    assert sections[0].bullets == ["x", "y"]


def test_parse_recovers_complete_sections_from_truncated_json() -> None:
    """A response cut off mid-document must still yield every COMPLETE section (not collapse to the
    tiny fallback). This is the core fix for half-completed documents."""
    from app.services.documents import DocumentService

    truncated = (
        '{"title": "Quarterly Plan", "subtitle": "FY26", "theme": "emerald", "sections": ['
        '{"heading": "Overview", "body": "Full first section.", "bullets": ["a", "b"]},'
        '{"heading": "Roadmap", "body": "Second section complete.", '
        '"table": {"headers": ["Phase","When"], "rows": [["1","Q1"],["2","Q2"]]}},'
        '{"heading": "Truncat'  # <-- the response was cut off here (no close)
    )
    parsed = DocumentService._parse(truncated)
    assert parsed is not None
    title, subtitle, theme, sections = parsed
    assert (title, subtitle, theme) == ("Quarterly Plan", "FY26", "emerald")
    assert [s.heading for s in sections] == ["Overview", "Roadmap"]  # the cut-off 3rd is dropped
    assert sections[1].table is not None and sections[1].table.headers == ["Phase", "When"]


def test_fallback_is_an_honest_notice_not_fabricated_content() -> None:
    """When the model returns nothing, the fallback must be an HONEST, on-topic notice — titled by
    the prompt, clearly stating content could not be generated, and with NO fabricated table/filler
    (the user complained the old scaffold served irrelevant boilerplate + a forced table)."""
    from app.services.documents import DocumentService

    title, subtitle, sections = DocumentService._fallback("Full documentation on how to use GitHub CI/CD")
    assert title.startswith("Full documentation on how to use GitHub CI/CD")  # on-topic title
    assert len(sections) == 1
    assert "could not" in sections[0].body.lower()  # plainly says it failed
    assert all(s.table is None for s in sections)  # never invents a table in the fallback


def test_latin1_transliterates_unicode_punctuation_for_pdf() -> None:
    """PDF latin-1 sanitization must transliterate smart punctuation to ASCII (not '?')."""
    from app.services.doc_render import _latin1

    out = _latin1("Scope creep — “quoted” café… bullet • arrow →")
    assert "—" not in out and "“" not in out and "•" not in out
    assert "?" not in out.replace("café", "")  # the accented é is latin-1-safe; no '?' from punctuation
    assert "-" in out and '"quoted"' in out and "..." in out and "->" in out


_RICH_SECTIONS = [
    {"heading": "Section 1", "body": "Body text.", "bullets": ["First point", "Second point ✓", "Third"]},
    {"heading": "Comparison", "body": "A small table.",
     "table": {"headers": ["Plan", "Price", "Notes"], "rows": [["Free", "$0", "Starter"], ["Pro", "$20", "Teams"]]}},
]


@pytest.mark.parametrize("fmt", ["pptx", "docx", "pdf"])
def test_real_render_produces_file(tmp_path, monkeypatch, fmt: str) -> None:
    """When the optional lib is installed, render a real, styled file (bullets + table + theme),
    bypassing the disable flag."""
    if not _engine_installed(fmt):
        pytest.skip(f"{_FMT_MODULE[fmt]} not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / f"doc.{fmt}"
    result = render_document(
        "Title ✓ unicode",  # exercises the latin-1 sanitizer for pdf
        _RICH_SECTIONS,
        out,
        fmt,
        subtitle="A styled subtitle ✓",
        theme="emerald",
    )
    assert result["ok"] and not result["stub"], result
    assert out.exists() and out.stat().st_size > 100


def test_structure_specs_cover_all_genres_and_are_distinct() -> None:
    """Every writing style maps to a StructureSpec, all 20 are present, and no two genres are
    byte-identical (so each renders as a visibly different document)."""
    from app.services.doc_render import STRUCTURE
    from app.services.documents import _STYLE_GUIDE

    assert set(STRUCTURE) == set(_STYLE_GUIDE)  # a structure for every offered genre, and vice-versa
    assert len(STRUCTURE) == 20
    assert len(set(STRUCTURE.values())) == 20, "two genres share an identical StructureSpec"
    # the high-impact axes (cover/columns/numbering/heading/container/bullet/divider) must vary a lot
    axes = {(s.cover, s.columns, s.numbering, s.heading_style, s.container, s.bullet, s.divider)
            for s in STRUCTURE.values()}
    assert len(axes) >= 16, "genres are too structurally similar"
    for name, s in STRUCTURE.items():
        assert s.image_max in (0, 1, 2), (name, s.image_max)
        assert s.image_policy in {"none", "optional", "inline", "hero", "fullbleed"}, (name, s.image_policy)
        if s.image_policy == "none":
            assert s.image_max == 0, name


def test_structure_spec_resolution_defaults() -> None:
    from app.services.doc_render import STRUCTURE, structure_spec

    assert structure_spec("LEGAL") is STRUCTURE["legal"]  # case-insensitive
    assert structure_spec("not-a-genre") is STRUCTURE["professional"]  # unknown -> default
    assert structure_spec(None) is STRUCTURE["professional"]


def test_theme_is_independent_of_style() -> None:
    """Color theme is USER-controlled, never derived from the genre: an explicit palette wins, 'auto'
    falls back to the agent suggestion then the default — the style never changes it."""
    from app.services.documents import _resolve_theme

    assert _resolve_theme("crimson", None) == "crimson"          # explicit wins
    assert _resolve_theme("auto", "forest") == "forest"          # auto -> agent suggestion
    assert _resolve_theme("auto", None) == "indigo"              # auto, no suggestion -> default
    assert _resolve_theme("nonsense", "teal") == "teal"          # bad request -> suggestion


def test_resolve_font_defaults() -> None:
    from app.services.documents import _resolve_font

    assert _resolve_font("SERIF") == "serif"
    assert _resolve_font("mono") == "mono"
    assert _resolve_font("comic") == "sans"  # unknown -> sans
    assert _resolve_font(None) == "sans"


@pytest.mark.parametrize(
    "font,fmt,pdf_core,office_font",
    [
        ("sans", "pdf", b"Helvetica", None),
        ("mono", "pdf", b"Courier", None),
        ("serif", "pdf", b"Times", None),
        ("sans", "docx", None, "Calibri"),
        ("mono", "docx", None, "Consolas"),
        ("serif", "docx", None, "Georgia"),
        ("sans", "pptx", None, "Calibri"),
        ("mono", "pptx", None, "Consolas"),
        ("serif", "pptx", None, "Georgia"),
    ],
)
def test_user_font_is_applied_independently(tmp_path, monkeypatch, font, fmt, pdf_core, office_font) -> None:
    """The USER's font choice (serif/sans/mono) lands in the output regardless of genre — Times/
    Helvetica/Courier for PDF; Georgia/Calibri/Consolas for DOCX & PPTX."""
    if not _engine_installed(fmt):
        pytest.skip(f"{_FMT_MODULE[fmt]} not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / f"doc.{fmt}"
    # 'legal' genre + every font: proves font is independent of the genre's own identity
    result = render_document("Styled Title", _RICH_SECTIONS, out, fmt, subtitle="sub", style="legal", font=font)
    assert result["ok"] and not result["stub"], result
    if fmt == "pdf":
        assert pdf_core in out.read_bytes()
    elif fmt == "docx":
        from docx import Document

        fonts = {r.font.name for para in Document(str(out)).paragraphs for r in para.runs if r.font.name}
        assert office_font in fonts, fonts
    else:  # pptx
        from pptx import Presentation

        fonts = {
            r.font.name
            for s in Presentation(str(out)).slides
            for sh in s.shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs
            for r in para.runs if r.font.name
        }
        assert office_font in fonts, fonts


@pytest.mark.parametrize("style", ["professional", "academic", "marketing", "legal", "technical", "creative", "friendly", "informative"])
def test_every_genre_renders_each_format(tmp_path, monkeypatch, style) -> None:
    """A representative spread of genres renders a real, non-empty file in all three formats."""
    for fmt in ("pdf", "docx", "pptx"):
        if not _engine_installed(fmt):
            continue
        monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
        out = tmp_path / f"{style}.{fmt}"
        result = render_document(f"{style} doc", _RICH_SECTIONS, out, fmt, subtitle="sub", style=style, theme="ocean")
        assert result["ok"] and not result["stub"], (style, fmt, result)
        assert out.exists() and out.stat().st_size > 500


def test_image_cap_enforced_and_skips_without_key(tmp_path, monkeypatch) -> None:
    """_attach_images caps generated images at spec.image_max, drops a text-only genre's requests,
    and (no HF key) leaves every image.path None so the document still renders."""
    import asyncio

    from app.schemas.documents import DocImage, DocSection
    from app.services.doc_render import structure_spec
    from app.services.documents import get_document_service

    svc = get_document_service()

    async def run(style, n_requested):
        secs = [DocSection(heading=f"S{i}", body="b", image=DocImage(prompt=f"img {i}"))
                for i in range(n_requested)]
        await svc._attach_images(secs, structure_spec(style), "Title", "doc_test", "default", [])
        return secs

    # text-only genre (legal): all image requests dropped
    legal = asyncio.get_event_loop().run_until_complete(run("legal", 3))
    assert all(s.image is None for s in legal)

    # image genre (casual, max 2): no HF key -> requests kept but capped to 2, paths all None
    casual = asyncio.get_event_loop().run_until_complete(run("casual", 4))
    kept = [s for s in casual if s.image is not None]
    assert len(kept) <= 2
    assert all(s.image.path is None for s in kept)


def test_section_image_embeds_when_path_present(tmp_path, monkeypatch) -> None:
    """A section image with a resolvable path is embedded by every renderer (the file gains the image)."""
    import struct
    import zlib

    img = tmp_path / "pic.png"
    w, h = 64, 48
    raw = bytearray()
    for _ in range(h):
        raw.append(0)
        raw += bytes((30, 110, 200)) * w

    def _chunk(typ, data):
        c = typ + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    img.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + _chunk(b"IDAT", zlib.compress(bytes(raw), 9))
        + _chunk(b"IEND", b"")
    )
    sections = [{"heading": "Hero", "body": "Body.", "image": {"path": "pic.png", "alt": "x"}}]
    for fmt in ("pdf", "docx", "pptx"):
        if not _engine_installed(fmt):
            continue
        monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
        out = tmp_path / f"img.{fmt}"
        r = render_document("Img", sections, out, fmt, subtitle="s", style="marketing", asset_root=tmp_path)
        assert r["ok"] and not r["stub"], (fmt, r)
        assert out.exists() and out.stat().st_size > 600


def test_unknown_theme_falls_back_without_error(tmp_path, monkeypatch) -> None:
    """An unknown theme name must resolve to the default palette, never raise."""
    if not _engine_installed("pdf"):
        pytest.skip("fpdf not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "doc.pdf"
    result = render_document("T", _RICH_SECTIONS, out, "pdf", theme="not-a-real-theme")
    assert result["ok"] and not result["stub"] and out.exists()


@pytest.mark.parametrize("fmt", ["pptx", "docx", "pdf"])
def test_degenerate_table_does_not_break_render(tmp_path, monkeypatch, fmt: str) -> None:
    """A malformed all-empty table ({headers:[], rows:[[],[]]}) must NOT crash any renderer
    (regression: it once raised ZeroDivisionError in PPTX -> whole deck degraded to markdown)."""
    if not _engine_installed(fmt):
        pytest.skip(f"{_FMT_MODULE[fmt]} not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / f"doc.{fmt}"
    sections = [{"heading": "Bad table", "body": "x", "table": {"headers": [], "rows": [[], []]}}]
    result = render_document("T", sections, out, fmt)
    assert result["ok"] and not result["stub"], result
    assert out.exists()


_CHART_SECTIONS = [
    {"heading": "Trend", "body": "Growth over the year.",
     "chart": {"type": "column", "title": "Revenue ($k)", "categories": ["Q1", "Q2", "Q3"],
               "series": [{"name": "2026", "values": [12, 19, 27]}]}},
    {"heading": "Split", "body": "Breakdown.",
     "chart": {"type": "pie", "title": "Share", "categories": ["A", "B"], "series": [{"name": "s", "values": [60, 40]}]}},
]


def test_pptx_renders_native_charts(tmp_path, monkeypatch) -> None:
    """A section with a chart must produce a native PPTX chart shape (a real graph, not a table)."""
    if not _engine_installed("pptx"):
        pytest.skip("pptx not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "deck.pptx"
    assert render_document("Charts", _CHART_SECTIONS, out, "pptx", theme="ocean")["ok"]
    from pptx import Presentation

    prs = Presentation(str(out))
    charts = sum(1 for s in prs.slides for sh in s.shapes if sh.has_chart)
    assert charts == 2, "expected a native chart on each chart-bearing slide"


@pytest.mark.parametrize("fmt", ["docx", "pdf"])
def test_chart_renders_as_table_in_documents(tmp_path, monkeypatch, fmt: str) -> None:
    """DOCX/PDF have no native chart engine — a chart must render as a captioned data table (no crash)."""
    if not _engine_installed(fmt):
        pytest.skip(f"{_FMT_MODULE[fmt]} not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / f"doc.{fmt}"
    assert render_document("Charts", _CHART_SECTIONS, out, fmt, theme="forest")["ok"]
    assert out.exists() and out.stat().st_size > 100
    if fmt == "docx":
        from docx import Document

        assert len(Document(str(out)).tables) == 2  # one data table per chart


def test_norm_chart_coerces_and_drops_empty() -> None:
    from app.services.doc_render import _norm_chart

    ct = _norm_chart({"type": "bar", "categories": ["a", "b"], "series": [{"name": "x", "values": ["1,200", "$3"]}]})
    assert ct is not None
    ctype, _title, cats, series = ct
    assert ctype == "bar" and cats == ["a", "b"] and series[0][1] == [1200.0, 3.0]
    assert _norm_chart({"series": []}) is None  # no numeric data -> not a chart
    assert _norm_chart("nope") is None


def test_num_rejects_non_finite() -> None:
    from app.services.doc_render import _num

    assert _num("inf") == 0.0 and _num("nan") == 0.0 and _num("1e400") == 0.0  # never inf/nan
    assert _num(float("inf")) == 0.0 and _num("12") == 12.0


def test_parse_chart_drops_all_non_numeric_and_non_finite() -> None:
    from app.services.documents import DocumentService

    assert DocumentService._parse_chart({"series": [{"name": "s", "values": ["abc", "--"]}]}) is None
    assert DocumentService._parse_chart({"series": [{"name": "s", "values": ["inf", "nan"]}]}) is None
    ok = DocumentService._parse_chart({"categories": ["a"], "series": [{"name": "s", "values": ["5", "x"]}]})
    assert ok is not None and ok.series[0].values == [5.0, 0.0]


def test_norm_chart_collapses_multi_series_pie() -> None:
    from app.services.doc_render import _norm_chart

    ct = _norm_chart({"type": "pie", "categories": ["a", "b"],
                      "series": [{"name": "one", "values": [1, 2]}, {"name": "two", "values": [3, 4]}]})
    assert ct is not None and len(ct[3]) == 1  # pie kept to a single series so all formats agree


def test_pptx_pie_uses_percentage_data_labels(tmp_path, monkeypatch) -> None:
    """A pie's data labels must show computed slice SHARE (show_percentage), not raw values as '7000%'."""
    if not _engine_installed("pptx"):
        pytest.skip("pptx not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "pie.pptx"
    sections = [{"heading": "Share", "chart": {"type": "pie", "title": "Split", "categories": ["A", "B", "C"],
                                               "series": [{"name": "s", "values": [10, 20, 70]}]}}]
    assert render_document("Pie", sections, out, "pptx", theme="rose")["ok"]
    from pptx import Presentation

    prs = Presentation(str(out))
    pie = next(sh.chart for s in prs.slides for sh in s.shapes if sh.has_chart)
    assert pie.plots[0].data_labels.show_percentage is True


def test_pptx_busy_slide_chart_stays_on_slide(tmp_path, monkeypatch) -> None:
    """A crowded slide (long body + many bullets + chart) must not push the chart off the 7.5in slide."""
    if not _engine_installed("pptx"):
        pytest.skip("pptx not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "busy.pptx"
    sections = [{
        "heading": "Crowded",
        "body": "A long body paragraph. " * 30,
        "bullets": [f"Bullet number {i} with a fair amount of descriptive text" for i in range(12)],
        "chart": {"type": "column", "title": "Data", "categories": ["a", "b", "c"], "series": [{"name": "s", "values": [1, 2, 3]}]},
    }]
    assert render_document("Busy deck", sections, out, "pptx", theme="indigo")["ok"]
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(out))
    slide_h = prs.slide_height
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_chart and sh.top is not None and sh.height is not None:
                assert sh.top + sh.height <= slide_h + Inches(0.1), "chart must stay within the slide"


def test_pptx_renders_both_chart_and_distinct_table(tmp_path, monkeypatch) -> None:
    """A section with BOTH a chart and a separate table keeps both in PPTX (parity with DOCX/PDF)."""
    if not _engine_installed("pptx"):
        pytest.skip("pptx not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "both.pptx"
    sections = [{
        "heading": "Both",
        "chart": {"type": "column", "categories": ["a", "b"], "series": [{"name": "s", "values": [1, 2]}]},
        "table": {"headers": ["K", "V"], "rows": [["x", "1"], ["y", "2"]]},
    }]
    assert render_document("Both", sections, out, "pptx", theme="teal")["ok"]
    from pptx import Presentation

    prs = Presentation(str(out))
    has_chart = any(sh.has_chart for s in prs.slides for sh in s.shapes)
    has_table = any(sh.has_table for s in prs.slides for sh in s.shapes)
    assert has_chart and has_table, "both the chart and the distinct table should render"


def test_pick_style_is_deterministic_and_in_set() -> None:
    from app.services.doc_render import STYLES, pick_style

    assert pick_style("GitHub CI/CD") == pick_style("GitHub CI/CD")  # stable per title
    assert pick_style("Quarterly Review") in STYLES


def test_parse_extracts_chart() -> None:
    from app.services.documents import DocumentService

    text = ('{"title":"T","sections":[{"heading":"H","chart":{"type":"line","title":"Trend",'
            '"categories":["Jan","Feb"],"series":[{"name":"users","values":[10,20]}]}}]}')
    parsed = DocumentService._parse(text)
    assert parsed is not None
    chart = parsed[3][0].chart
    assert chart is not None and chart.type == "line" and chart.series[0].values == [10.0, 20.0]


def test_pptx_has_entrance_animations(tmp_path, monkeypatch) -> None:
    """Content-bearing slides must carry an injected <p:timing> entrance animation tree."""
    if not _engine_installed("pptx"):
        pytest.skip("pptx not installed")
    monkeypatch.delenv("OMNIVRA_DISABLE_RENDER", raising=False)
    out = tmp_path / "deck.pptx"
    assert render_document("Deck", _RICH_SECTIONS, out, "pptx", subtitle="Sub", theme="violet")["ok"]

    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(str(out))
    animated = [s for s in prs.slides if s._element.find(qn("p:timing")) is not None]
    assert animated, "expected at least one slide with an entrance-animation <p:timing> tree"
